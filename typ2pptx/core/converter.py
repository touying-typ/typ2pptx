"""
Main converter module: typst.ts SVG → PPTX.

This module orchestrates the conversion from typst.ts SVG to PowerPoint,
using the ppt-master svg_to_shapes.py as the foundation and extending it
for typst.ts specific SVG structure.
"""

import os
import re
import math
import json
import copy
import shutil
import tempfile
import subprocess
from pathlib import Path
from typing import Optional, Tuple, List, Dict, Any, Set
from xml.etree import ElementTree as ET
from dataclasses import dataclass, field
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Cm, Pt, Emu
from pptx.enum.text import PP_ALIGN

from .typst_svg_parser import (
    parse_typst_svg, TypstSVGData, PageData, TextSegment,
    ShapeElement, GlyphInfo, GlyphUse, FontVariant, LinkRegion,
    parse_transform, SVG_NS, XLINK_NS,
)


# EMU per SVG pixel (at 96 DPI)
EMU_PER_PX = 9525


@dataclass
class ConversionConfig:
    """Configuration for the SVG to PPTX conversion."""
    # Whether to use the adapted svg_to_shapes for shape conversion
    use_native_shapes: bool = True
    # DPI for rasterization fallback
    raster_dpi: int = 300
    # Default font for text
    default_latin_font: str = "Arial"
    default_ea_font: str = "Microsoft YaHei"
    default_cs_font: str = "Arial"
    # Whether to include speaker notes
    include_speaker_notes: bool = True
    # Verbose output
    verbose: bool = False
    # Math rendering mode for inline math: "text", "glyph", or "auto" (heuristic)
    # "auto" uses text for simple formulas (letters, digits, sub/superscripts) and
    # glyph curves for complex ones (integrals, matrices, large operators)
    inline_math_mode: str = "auto"
    # Math rendering mode for display/block math: "text", "glyph", or "auto"
    # "auto" uses the same heuristic as inline math
    display_math_mode: str = "glyph"


def _compute_text_position(text_seg: TextSegment, page_width: float, page_height: float) -> Tuple[float, float, float, float]:
    """Compute the position and size of a text segment in SVG pixels.

    Returns (x, y, width, height) in SVG pixels.
    """
    x = text_seg.x
    y = text_seg.y
    width = text_seg.width
    height = text_seg.height

    return (x, y, width, height)


def _font_variant_to_props(variant: str) -> Dict[str, Any]:
    """Convert font variant name to PowerPoint text properties."""
    props = {
        'bold': False,
        'italic': False,
        'font_family': 'Arial',
    }

    if variant in ('bold', 'bolditalic'):
        props['bold'] = True
    if variant in ('italic', 'bolditalic'):
        props['italic'] = True
    if variant == 'mono':
        props['font_family'] = 'Consolas'
    if variant == 'math':
        props['font_family'] = 'Cambria Math'

    return props


_NAMED_COLORS = {
    "black": "000000", "white": "FFFFFF", "red": "FF0000", "green": "008000",
    "blue": "0000FF", "yellow": "FFFF00", "cyan": "00FFFF", "magenta": "FF00FF",
    "orange": "FFA500", "purple": "800080", "pink": "FFC0CB", "brown": "A52A2A",
    "gray": "808080", "grey": "808080", "silver": "C0C0C0", "gold": "FFD700",
    "navy": "000080", "teal": "008080", "maroon": "800000", "olive": "808000",
    "lime": "00FF00", "aqua": "00FFFF", "fuchsia": "FF00FF",
    "darkgray": "A9A9A9", "darkgrey": "A9A9A9", "lightgray": "D3D3D3",
    "lightgrey": "D3D3D3", "darkred": "8B0000", "darkgreen": "006400",
    "darkblue": "00008B", "coral": "FF7F50", "salmon": "FA8072",
    "tomato": "FF6347", "crimson": "DC143C", "indigo": "4B0082",
    "violet": "EE82EE", "turquoise": "40E0D0", "tan": "D2B48C",
    "khaki": "F0E68C", "ivory": "FFFFF0", "beige": "F5F5DC",
    "linen": "FAF0E6", "lavender": "E6E6FA", "plum": "DDA0DD",
    "orchid": "DA70D6", "peru": "CD853F", "sienna": "A0522D",
    "chocolate": "D2691E", "wheat": "F5DEB3", "snow": "FFFAFA",
    "honeydew": "F0FFF0", "azure": "F0FFFF", "steelblue": "4682B4",
    "royalblue": "4169E1", "dodgerblue": "1E90FF", "skyblue": "87CEEB",
    "slategray": "708090", "slategrey": "708090",
}


def _parse_color(color_str: str) -> Optional[str]:
    """Parse color string to RRGGBB hex (discards alpha).

    Supports: #RGB, #RRGGBB, #RRGGBBAA, rgb(r,g,b), rgba(r,g,b,a), named CSS colors.
    Returns None for 'none', empty, or url() references.
    """
    result = _parse_color_with_alpha(color_str)
    if result is None:
        return None
    return result[0]


def _parse_color_with_alpha(color_str: str) -> Optional[Tuple[str, float]]:
    """Parse color string to (RRGGBB hex, alpha float 0.0-1.0).

    Supports: #RGB, #RRGGBB, #RRGGBBAA, rgb(r,g,b), rgba(r,g,b,a), named CSS colors.
    Returns None for 'none', empty, or url() references.
    Alpha of 1.0 means fully opaque.
    """
    if not color_str:
        return None
    color_str = color_str.strip()

    # Skip gradients and other url references
    if color_str.startswith('url(') or color_str == 'none':
        return None

    # Hex colors
    if color_str.startswith('#'):
        hex_val = color_str[1:]
        if len(hex_val) == 3:
            hex_val = ''.join(c * 2 for c in hex_val)
        if len(hex_val) == 6:
            return (hex_val.upper(), 1.0)
        # 8-digit hex (with alpha) - RRGGBBAA format
        if len(hex_val) == 8:
            rgb_hex = hex_val[:6].upper()
            alpha = int(hex_val[6:8], 16) / 255.0
            return (rgb_hex, alpha)
        return None

    # rgb() / rgba()
    m = re.match(r'rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*(?:,\s*([\d.]+)\s*)?\)', color_str)
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        alpha = float(m.group(4)) if m.group(4) else 1.0
        if alpha > 1.0:
            alpha = alpha / 255.0  # Typst uses 0-255 range for alpha
        return (f'{r:02X}{g:02X}{b:02X}', alpha)

    # rgb with percentages
    m = re.match(r'rgba?\(\s*([\d.]+)%\s*,\s*([\d.]+)%\s*,\s*([\d.]+)%\s*(?:,\s*([\d.]+)%?\s*)?\)', color_str)
    if m:
        r = int(float(m.group(1)) * 255 / 100)
        g = int(float(m.group(2)) * 255 / 100)
        b = int(float(m.group(3)) * 255 / 100)
        alpha = float(m.group(4)) / 100.0 if m.group(4) else 1.0
        return (f'{min(r,255):02X}{min(g,255):02X}{min(b,255):02X}', alpha)

    # Named colors
    lower = color_str.lower()
    if lower in _NAMED_COLORS:
        return (_NAMED_COLORS[lower], 1.0)

    return None


def _apply_alpha_to_shape(shape, alpha: float):
    """Apply alpha transparency to a python-pptx shape's fill.

    Args:
        shape: A python-pptx shape with fill.solid() already called
        alpha: Alpha value 0.0 (transparent) to 1.0 (opaque)
    """
    if alpha >= 1.0:
        return
    from lxml import etree
    # Access the fill element and add alpha child
    fill_elem = shape.fill._fill
    if fill_elem is not None:
        # Find the srgbClr or any color element
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        for color_elem in fill_elem.iter():
            tag = color_elem.tag.split('}')[-1] if '}' in color_elem.tag else color_elem.tag
            if tag in ('srgbClr', 'schemeClr', 'prstClr'):
                alpha_val = int(alpha * 100000)
                alpha_elem = etree.SubElement(
                    color_elem,
                    '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'
                )
                alpha_elem.set('val', str(alpha_val))
                return


def _split_line_by_gaps(line_segs: list, gap_threshold_factor: float = 1.0) -> list:
    """Split a line group into sub-groups when there are large horizontal gaps.

    This handles cases like footer text where "Authors", "Title", "date",
    and page numbers are on the same baseline but spread far apart. These
    should be separate textboxes, not one merged textbox.

    Args:
        line_segs: List of segments on the same line, sorted by x position.
        gap_threshold_factor: A gap larger than this × avg_font_size triggers a split.

    Returns:
        List of segment lists (sub-groups).
    """
    if len(line_segs) <= 1:
        return [line_segs]

    sorted_segs = sorted(line_segs, key=lambda s: s.x)
    avg_font_size = sum(s.font_size for s in sorted_segs) / len(sorted_segs)
    gap_threshold = avg_font_size * gap_threshold_factor

    sub_groups = []
    current = [sorted_segs[0]]

    for i in range(1, len(sorted_segs)):
        prev = sorted_segs[i - 1]
        cur = sorted_segs[i]
        gap = cur.x - (prev.x + prev.width)

        if gap > gap_threshold:
            sub_groups.append(current)
            current = [cur]
        else:
            current.append(cur)

    sub_groups.append(current)
    return sub_groups


def _consolidate_dot_leaders(segments: list) -> list:
    """Consolidate individual dot leader characters into single segments.

    TOC/outline pages often have many individual '.' segments used as dot
    leaders between an entry title and its page number. These create ugly
    textboxes with dozens of individual period characters.

    This function groups consecutive single-dot segments on the same baseline
    into a single consolidated segment, making the output much cleaner.
    """
    if not segments:
        return segments

    # Characters that can be dot leaders
    DOT_CHARS = set('.·…⁠')  # includes zero-width space U+2060

    def get_baseline(seg):
        return seg.y + seg.font_size * 0.8

    # Sort by baseline then x position
    sorted_segs = sorted(segments, key=lambda s: (round(get_baseline(s), 1), s.x))

    result = []
    i = 0
    while i < len(sorted_segs):
        seg = sorted_segs[i]
        text = seg.text.strip()

        # Check if this is a single dot character
        if text in DOT_CHARS:
            # Collect consecutive dots on the same baseline
            dot_group = [seg]
            j = i + 1
            while j < len(sorted_segs):
                next_seg = sorted_segs[j]
                next_text = next_seg.text.strip()
                # Same baseline (within tolerance)?
                if abs(get_baseline(next_seg) - get_baseline(seg)) > 2.0:
                    break
                # Is it another dot?
                if next_text not in DOT_CHARS:
                    break
                # Is it close enough horizontally? (within 2× font size)
                prev_right = dot_group[-1].x + dot_group[-1].width
                if next_seg.x - prev_right > seg.font_size * 2:
                    break
                dot_group.append(next_seg)
                j += 1

            if len(dot_group) >= 3:
                # Consolidate: create a single segment with combined dots
                # Use first segment as template, extend its width
                import copy
                consolidated = copy.copy(dot_group[0])
                consolidated.text = '.' * len(dot_group)
                last = dot_group[-1]
                consolidated.width = (last.x + last.width) - consolidated.x
                result.append(consolidated)
                i = j
            else:
                # Not enough dots to consolidate, keep as-is
                result.append(seg)
                i += 1
        else:
            result.append(seg)
            i += 1

    return result


def _group_segments_by_line(segments: list) -> list:
    """Group text segments that are on the same line.

    Two segments are on the same line if they have approximately the same
    baseline Y position (within tolerance). Font size differences are allowed
    since inline code, links, etc. may have different sizes on the same line.

    After grouping by baseline, each line group is further split at large
    horizontal gaps (e.g., footer elements spread across the slide).

    Returns a list of lists, where each inner list contains segments
    from the same line, sorted by baseline Y.
    """
    if not segments:
        return []

    # Reconstruct baseline_y from the stored top_y:
    # top_y = baseline_y - font_size * 0.8
    # So: baseline_y = top_y + font_size * 0.8
    def get_baseline(seg):
        return seg.y + seg.font_size * 0.8

    # Sort by baseline_y
    sorted_segs = sorted(segments, key=lambda s: get_baseline(s))

    baseline_groups = []
    current_group = [sorted_segs[0]]
    current_baseline = get_baseline(sorted_segs[0])

    for seg in sorted_segs[1:]:
        seg_baseline = get_baseline(seg)

        # Tolerance: baselines within 2px are considered the same line
        if abs(seg_baseline - current_baseline) < 2.0:
            current_group.append(seg)
        else:
            baseline_groups.append(current_group)
            current_group = [seg]
            current_baseline = seg_baseline

    baseline_groups.append(current_group)

    # Split each baseline group at large horizontal gaps
    final_groups = []
    for group in baseline_groups:
        sub_groups = _split_line_by_gaps(group)
        final_groups.extend(sub_groups)

    return final_groups


def _merge_inline_math_subscripts(line_groups: list) -> list:
    """Merge inline math sub/superscript line groups into their adjacent text lines.

    When inline math like $A = pi r^2$ is grouped by baseline, the superscript '2'
    ends up in a separate line group because its baseline differs from the main line
    by ~9px (well beyond the 2px tolerance). This function identifies such isolated
    small-font math segments and merges them into the nearest horizontally-overlapping
    text line group.

    Criteria for a "sub/superscript" line group:
    - All segments in the group have font_variant == 'math'
    - The group's font size is smaller than the dominant font size of nearby lines

    Merge target: the nearest line group (by baseline distance) that horizontally
    overlaps and has larger-font segments.
    """
    if len(line_groups) <= 1:
        return line_groups

    def get_baseline(seg):
        return seg.y + seg.font_size * 0.8

    def group_baseline(group):
        return sum(get_baseline(s) for s in group) / len(group)

    def group_x_range(group):
        return (min(s.x for s in group), max(s.x + s.width for s in group))

    def group_max_font_size(group):
        return max(s.font_size for s in group)

    # Identify which groups are "isolated math subscript/superscript" groups
    # vs "main text" groups
    is_math_sub = [False] * len(line_groups)
    for i, group in enumerate(line_groups):
        # Check if all segments are math with small font
        if all(s.font_variant == 'math' for s in group):
            avg_size = sum(s.font_size for s in group) / len(group)
            # Check if there's a nearby group with larger font
            for j, other_group in enumerate(line_groups):
                if i == j:
                    continue
                other_max_size = group_max_font_size(other_group)
                if avg_size < other_max_size * 0.85:
                    # Check horizontal overlap with the other group
                    my_xmin, my_xmax = group_x_range(group)
                    other_xmin, other_xmax = group_x_range(other_group)
                    # Segments should be within font-size distance horizontally
                    x_gap = max(0, max(my_xmin - other_xmax, other_xmin - my_xmax))
                    if x_gap < other_max_size * 2:
                        # Check baseline proximity (sub/super are within ~1.5x font size)
                        bl_diff = abs(group_baseline(group) - group_baseline(other_group))
                        if bl_diff < other_max_size * 0.8:
                            is_math_sub[i] = True
                            break

    if not any(is_math_sub):
        return line_groups

    # Build a new list: merge math sub groups into their best target
    result = []
    merged_into = {}  # index -> target index

    for i, group in enumerate(line_groups):
        if not is_math_sub[i]:
            continue
        # Find the best merge target
        best_target = None
        best_bl_diff = float('inf')
        my_bl = group_baseline(group)
        my_xmin, my_xmax = group_x_range(group)

        for j, other_group in enumerate(line_groups):
            if i == j or is_math_sub[j]:
                continue
            other_bl = group_baseline(other_group)
            bl_diff = abs(my_bl - other_bl)
            other_max_size = group_max_font_size(other_group)
            if bl_diff < other_max_size * 0.8:
                other_xmin, other_xmax = group_x_range(other_group)
                x_gap = max(0, max(my_xmin - other_xmax, other_xmin - my_xmax))
                if x_gap < other_max_size * 2:
                    if bl_diff < best_bl_diff:
                        best_bl_diff = bl_diff
                        best_target = j

        if best_target is not None:
            merged_into[i] = best_target

    # Build result
    merge_additions = {}  # target_index -> list of segments to add
    for src, tgt in merged_into.items():
        if tgt not in merge_additions:
            merge_additions[tgt] = []
        merge_additions[tgt].extend(line_groups[src])

    for i, group in enumerate(line_groups):
        if i in merged_into:
            continue  # This group was merged into another
        if i in merge_additions:
            result.append(group + merge_additions[i])
        else:
            result.append(group)

    return result


def _is_simple_inline_math(cluster: list) -> bool:
    """Determine if an inline math cluster is simple enough for text rendering.

    Simple inline math contains only:
    - Normal letters, digits, basic operators (+, -, =, <, >, etc.)
    - Greek letters (single Unicode chars)
    - Sub/superscripts (detected by font size variation)
    - Parentheses, brackets

    Complex inline math (returns False) contains:
    - Too many distinct font sizes (suggests fractions or nested structures)
    - Very small font sizes relative to the dominant (suggests stacked fractions)
    - Vertically stacked segments at the same x (stacked fractions like frac(1,2))
    - Complex math symbols (integrals, sums, roots, etc.)
    """
    if not cluster:
        return True

    # Get font size statistics
    sizes = [s.font_size for s in cluster]
    dominant_size = max(sizes)
    min_size = min(sizes)

    # If the smallest character is less than 60% of the dominant size,
    # it's likely a fraction or complex nested structure
    if min_size < dominant_size * 0.55:
        return False

    # Count distinct font size levels
    # Simple math has at most 2 levels (normal + sub/superscript)
    size_levels = set()
    for s in sizes:
        # Round to nearest 0.5 to group similar sizes
        size_levels.add(round(s * 2) / 2)
    if len(size_levels) > 3:
        return False

    # Detect vertically stacked segments at similar x positions
    # This catches stacked fractions like frac(1,2) where numerator and
    # denominator are at the same x but different y, both at reduced size.
    # Important: sub/superscript pairs (like a_1^2) also stack vertically,
    # but they always have a normal-sized character (like 'a') at a nearby x.
    # Fractions have the stacked pair isolated from normal-sized characters.
    if len(cluster) >= 2:
        small_segs = [s for s in cluster if s.font_size < dominant_size * 0.85]
        normal_segs = [s for s in cluster if s.font_size >= dominant_size * 0.85]
        if len(small_segs) >= 2:
            for i, a in enumerate(small_segs):
                for b in small_segs[i + 1:]:
                    # Similar x position (within half the segment width)
                    x_overlap = min(a.x + a.width, b.x + b.width) - max(a.x, b.x)
                    min_w = min(a.width, b.width)
                    if min_w > 0 and x_overlap > min_w * 0.3:
                        # Vertically separated
                        a_bl = a.y + a.font_size * 0.8
                        b_bl = b.y + b.font_size * 0.8
                        bl_diff = abs(a_bl - b_bl)
                        if bl_diff > dominant_size * 0.3:
                            # Check if there's a normal-sized segment right next to
                            # this stacked pair (indicating sub/superscript, not fraction)
                            pair_x_min = min(a.x, b.x)
                            pair_x_max = max(a.x + a.width, b.x + b.width)
                            has_adjacent_normal = False
                            for ns in normal_segs:
                                ns_right = ns.x + ns.width
                                # Adjacent if gap is small (within 2px)
                                if (abs(ns_right - pair_x_min) < 2 or
                                        abs(pair_x_max - ns.x) < 2):
                                    has_adjacent_normal = True
                                    break
                            if not has_adjacent_normal:
                                return False  # Stacked fraction detected

    # Check character content - look for complex math symbols
    # that don't render well as text
    complex_chars = set('∫∑∏√∛∜⋂⋃⨁⨂')
    for seg in cluster:
        text = seg.text.strip()
        for ch in text:
            if ch in complex_chars:
                return False

    # If there are many segments (>15), it's likely complex
    if len(cluster) > 15:
        return False

    return True


def _merge_paragraph_lines(line_groups: list, page_width: float = 0) -> list:
    """Merge consecutive single-line groups that form a paragraph.

    Detects consecutive lines with:
    - Similar left x position (within tolerance)
    - Same font properties (variant, size, color)
    - Consistent vertical spacing between lines
    - Each line is a single text segment (i.e., wrapped text, not multi-run)
      OR multiple segments that share the same properties

    These are merged into a single "paragraph group" that will be rendered
    as one textbox with word wrapping enabled.

    Supports multi-column layouts (e.g., Typst's #columns(2)[...]) by first
    clustering line groups into column regions based on their x positions,
    then merging paragraphs within each column independently using the
    column's actual width rather than the full page width.

    Returns a list of dicts:
    - {'type': 'line', 'segments': [...]}  for regular single-line groups
    - {'type': 'paragraph', 'lines': [[...], [...], ...], 'segments': [...]}
      for merged paragraph groups
    """
    if not line_groups:
        return []

    def _line_props(line_segs):
        """Get the dominant properties of a line group."""
        if not line_segs:
            return None
        seg = line_segs[0]
        return {
            'font_variant': seg.font_variant,
            'font_size': round(seg.font_size, 1),
            'fill_color': seg.fill_color,
        }

    def _line_left_x(line_segs):
        """Get the leftmost x of a line group."""
        return min(s.x for s in line_segs)

    def _line_baseline_y(line_segs):
        """Get the average baseline y of a line group."""
        return sum(s.y + s.font_size * 0.8 for s in line_segs) / len(line_segs)

    _BULLET_CHARS = set('•‣⁃◦▸▹►▻–—-∗⋅·')

    def _line_starts_with_bullet(line_segs):
        """Check if a line starts with a bullet/list marker character."""
        if not line_segs:
            return False
        sorted_by_x = sorted(line_segs, key=lambda s: s.x)
        first_text = sorted_by_x[0].text.strip()
        if first_text and first_text[0] in _BULLET_CHARS:
            return True
        return False

    def _line_text_width(line_segs):
        """Get the total text width span of a line."""
        if not line_segs:
            return 0
        min_x = min(s.x for s in line_segs)
        max_right = max(s.x + s.width for s in line_segs)
        return max_right - min_x

    def _line_right_x(line_segs):
        """Get the rightmost x coordinate of a line group."""
        if not line_segs:
            return 0
        return max(s.x + s.width for s in line_segs)

    _DOT_LEADER_CHARS = set('.·…⁠')  # includes zero-width space U+2060

    def _line_has_dot_leaders(line_segs):
        """Check if a line contains dot leader patterns (TOC entries)."""
        for seg in line_segs:
            text = seg.text.strip()
            dot_count = 0
            for ch in text:
                if ch in _DOT_LEADER_CHARS:
                    dot_count += 1
                    if dot_count >= 3:
                        return True
                else:
                    dot_count = 0
        return False

    def _line_has_numbering(line_segs):
        """Check if a line starts with a numbering pattern (e.g. '1.', '1.1', etc)."""
        if not line_segs:
            return False
        sorted_by_x = sorted(line_segs, key=lambda s: s.x)
        first_text = sorted_by_x[0].text.strip()
        import re
        if re.match(r'^\d+\.', first_text):
            return True
        return False

    def _line_dominant_font_size(line_segs):
        """Get the dominant (most common) font size of a line group."""
        if not line_segs:
            return 0
        sizes = [round(s.font_size, 1) for s in line_segs]
        from collections import Counter
        return Counter(sizes).most_common(1)[0][0]

    def _cluster_into_columns(line_groups_list):
        """Cluster line groups into column regions based on x positions.

        Groups lines whose left x positions are within 5px of each other
        into the same column cluster. Returns a list of (column_lines, column_info)
        tuples, where column_info contains the effective width for that column.

        Each column's lines are sorted by baseline Y for proper paragraph detection.
        """
        if not line_groups_list:
            return []

        # Collect (left_x, original_index, line_group) for clustering
        indexed_lines = []
        for idx, lg in enumerate(line_groups_list):
            left_x = _line_left_x(lg)
            indexed_lines.append((left_x, idx, lg))

        # Sort by left_x to find column clusters
        indexed_lines.sort(key=lambda t: t[0])

        # Cluster by left_x proximity (within 5px)
        columns = []
        current_cluster = [indexed_lines[0]]
        current_x = indexed_lines[0][0]

        for item in indexed_lines[1:]:
            if abs(item[0] - current_x) <= 5.0:
                current_cluster.append(item)
            else:
                columns.append(current_cluster)
                current_cluster = [item]
                current_x = item[0]
        columns.append(current_cluster)

        # For each column, compute effective width and sort by baseline Y
        result = []
        for cluster in columns:
            # Sort lines within this column by baseline Y
            cluster.sort(key=lambda t: _line_baseline_y(t[2]))
            column_lines = [t[2] for t in cluster]
            original_indices = [t[1] for t in cluster]

            # Compute the effective content width for this column:
            # the maximum right_x minus the minimum left_x across all lines
            col_left = min(_line_left_x(lg) for lg in column_lines)
            col_right = max(_line_right_x(lg) for lg in column_lines)
            col_width = col_right - col_left

            result.append({
                'lines': column_lines,
                'indices': original_indices,
                'left': col_left,
                'right': col_right,
                'width': col_width,
            })

        return result

    def _can_merge(line_a, line_b, effective_width, effective_right,
                   expected_spacing=None):
        """Check if two consecutive lines can be merged into a paragraph.

        Args:
            line_a: First line group.
            line_b: Second line group (immediately below line_a in the same column).
            effective_width: The content width of the column/region these lines
                belong to, used for width threshold checks instead of page_width.
            effective_right: The right boundary of the column/region.
            expected_spacing: Expected vertical spacing from previous merges.
        """
        props_a = _line_props(line_a)
        props_b = _line_props(line_b)

        if props_a is None or props_b is None:
            return False, None

        if _line_starts_with_bullet(line_a) or _line_starts_with_bullet(line_b):
            return False, None

        if _line_has_dot_leaders(line_a) or _line_has_dot_leaders(line_b):
            return False, None

        if _line_has_numbering(line_b):
            return False, None

        dominant_size_a = _line_dominant_font_size(line_a)
        dominant_size_b = _line_dominant_font_size(line_b)
        if dominant_size_a != dominant_size_b:
            return False, None

        x_a = _line_left_x(line_a)
        x_b = _line_left_x(line_b)
        if abs(x_a - x_b) > 5.0:
            return False, None

        # For paragraph wrapping, line_a should be wide enough.
        # Use effective_width (column width) instead of page_width so that
        # multi-column layouts are handled correctly.
        width_a = _line_text_width(line_a)
        reference_width = effective_width if effective_width > 0 else (
            page_width if page_width > 0 else 800
        )
        min_paragraph_width = reference_width * 0.44
        if width_a < min_paragraph_width:
            return False, None

        # Explicit line break detection: line_a should reach near the right
        # boundary of its column/region.
        right_a = _line_right_x(line_a)
        right_boundary = effective_right if effective_right > 0 else (
            page_width if page_width > 0 else 800
        )
        right_threshold = right_boundary * 0.70
        # For columns that don't start at x=0, use column-relative check:
        # line_a should fill at least 70% of the column's width
        if effective_width > 0:
            col_left = _line_left_x(line_a)
            line_fill_ratio = width_a / effective_width
            if line_fill_ratio < 0.70:
                return False, None
        elif right_a < right_threshold:
            return False, None

        # Compute vertical spacing
        bl_a = _line_baseline_y(line_a)
        bl_b = _line_baseline_y(line_b)
        spacing = bl_b - bl_a

        font_size = dominant_size_a
        if spacing <= 0 or spacing > font_size * 2.0:
            return False, None

        if expected_spacing is not None:
            if abs(spacing - expected_spacing) > expected_spacing * 0.20:
                return False, None

        return True, spacing

    def _merge_column_lines(column_lines, effective_width, effective_right):
        """Merge paragraph lines within a single column."""
        column_result = []
        i = 0
        while i < len(column_lines):
            para_lines = [column_lines[i]]
            spacings = []

            j = i + 1
            while j < len(column_lines):
                expected = spacings[-1] if spacings else None
                can_do_merge, spacing = _can_merge(
                    para_lines[-1], column_lines[j],
                    effective_width, effective_right, expected
                )
                if can_do_merge:
                    para_lines.append(column_lines[j])
                    spacings.append(spacing)
                    j += 1
                else:
                    break

            if len(para_lines) >= 2:
                all_segs = []
                for line in para_lines:
                    all_segs.extend(line)
                column_result.append({
                    'type': 'paragraph',
                    'lines': para_lines,
                    'segments': all_segs,
                    'line_spacing': spacings[0] if spacings else 0,
                })
            else:
                for line in para_lines:
                    column_result.append({
                        'type': 'line',
                        'segments': line,
                    })

            i = j
        return column_result

    # Detect column layout: cluster line groups by x position
    columns = _cluster_into_columns(line_groups)

    if len(columns) <= 1:
        # Single column: use page_width as before (backward compatible)
        effective_width = page_width if page_width > 0 else 0
        effective_right = page_width if page_width > 0 else 0
        sorted_lines = columns[0]['lines'] if columns else line_groups
        return _merge_column_lines(sorted_lines, effective_width, effective_right)

    # Multi-column layout: merge each column independently, then interleave
    # results back in original baseline Y order
    all_results_with_y = []
    for col_info in columns:
        col_lines = col_info['lines']
        col_width = col_info['width']
        col_right = col_info['right']
        merged = _merge_column_lines(col_lines, col_width, col_right)
        for group in merged:
            # Compute a representative Y for ordering
            segs = group['segments'] if group['type'] == 'line' else group['segments']
            rep_y = min(s.y for s in segs)
            rep_x = min(s.x for s in segs)
            all_results_with_y.append((rep_y, rep_x, group))

    # Sort by Y first, then X to maintain visual order
    all_results_with_y.sort(key=lambda t: (t[0], t[1]))
    return [item[2] for item in all_results_with_y]


def _detect_font_style_by_prefix(
    prefix: str,
    font_variants: Dict[str, FontVariant],
    prefix_to_style: Dict[str, str],
) -> str:
    """Detect font style from glyph prefix."""
    if prefix in prefix_to_style:
        return prefix_to_style[prefix]
    if prefix in font_variants:
        return font_variants[prefix].style
    return "regular"


def _analyze_font_variants_from_svg(svg_data: TypstSVGData) -> Dict[str, str]:
    """Analyze the SVG data to map glyph prefixes to font styles.

    Strategy:
    1. The most used prefix is 'regular'
    2. Prefixes using quadratic curves are 'mono'
    3. Others need context analysis from text groups

    Returns dict mapping prefix -> style
    """
    prefix_to_style: Dict[str, str] = {}

    # Sort by usage count (most used = regular)
    sorted_variants = sorted(
        svg_data.font_variants.items(),
        key=lambda x: x[1].glyph_count,
        reverse=True,
    )

    if not sorted_variants:
        return prefix_to_style

    # First, mark quadratic-curve prefixes as 'mono'
    for prefix, variant in sorted_variants:
        if variant.uses_quadratic:
            prefix_to_style[prefix] = 'mono'

    # The most-used non-mono prefix is 'regular'
    for prefix, variant in sorted_variants:
        if prefix not in prefix_to_style:
            prefix_to_style[prefix] = 'regular'
            break

    # Remaining prefixes need more analysis - we'll try to detect from the
    # text segments by looking at fill colors and context
    # For now, label them by usage count:
    # - Second most used (non-mono, non-regular) → could be bold (headings reuse)
    # - Less used → italic, bolditalic
    remaining = [p for p, v in sorted_variants if p not in prefix_to_style]

    # We need to look at the actual text content to distinguish
    # bold vs italic vs bolditalic. For now, use a heuristic:
    # We'll refine this during text processing
    for prefix in remaining:
        prefix_to_style[prefix] = 'unknown'  # Will be refined

    return prefix_to_style


class TypstSVGConverter:
    """Converts typst.ts SVG pages to PPTX slides."""

    def __init__(self, config: ConversionConfig = None):
        self.config = config or ConversionConfig()
        self._svg_data: Optional[TypstSVGData] = None

    def convert(self, svg_path: str, output_path: str,
                speaker_notes: Optional[Dict[int, str]] = None) -> str:
        """Convert a typst.ts SVG file to PPTX.

        Args:
            svg_path: Path to the SVG file from typst-ts-cli
            output_path: Output PPTX file path
            speaker_notes: Optional dict mapping page_num (0-indexed) to notes text

        Returns:
            Path to the generated PPTX file
        """
        # Parse the SVG
        self._svg_data = parse_typst_svg(svg_path)

        if self.config.verbose:
            print(f"Parsed SVG: {len(self._svg_data.pages)} pages, "
                  f"viewBox={self._svg_data.viewbox_width}x{self._svg_data.viewbox_height}")
            print(f"Font variants: {len(self._svg_data.font_variants)}")
            for prefix, variant in self._svg_data.font_variants.items():
                print(f"  {prefix}: {variant.glyph_count} glyphs, "
                      f"quadratic={variant.uses_quadratic}")

        # Use the font variant styles already computed by the parser's prescan
        prefix_to_style = {
            prefix: variant.style
            for prefix, variant in self._svg_data.font_variants.items()
        }

        if self.config.verbose:
            print(f"Font style mapping: {prefix_to_style}")

        # Create presentation
        prs = Presentation()

        # Set slide dimensions based on SVG viewBox
        page_width = self._svg_data.viewbox_width
        if self._svg_data.pages:
            page_height = self._svg_data.pages[0].height
        else:
            page_height = self._svg_data.viewbox_height

        # Set slide size in EMU
        prs.slide_width = Emu(int(page_width * EMU_PER_PX))
        prs.slide_height = Emu(int(page_height * EMU_PER_PX))

        if self.config.verbose:
            print(f"Slide size: {page_width}x{page_height} px = "
                  f"{prs.slide_width}x{prs.slide_height} EMU")

        blank_layout = prs.slide_layouts[6]  # Blank layout

        # Convert each page
        for page_data in self._svg_data.pages:
            slide = prs.slides.add_slide(blank_layout)

            # Convert shapes using ppt-master's svg_to_shapes approach
            self._convert_page_shapes(
                page_data, slide, page_width, page_height,
                self._svg_data, prefix_to_style,
            )

            # Convert text segments
            self._convert_page_texts(
                page_data, slide, page_width, page_height,
                prefix_to_style,
            )

            # Apply hyperlinks from SVG <a> elements
            if page_data.links:
                self._apply_links(page_data, slide)

            # Add speaker notes if available
            if speaker_notes and (page_data.page_num - 1) in speaker_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = speaker_notes[page_data.page_num - 1]

        # Override theme hyperlink colors to prevent PowerPoint's default
        # blue (hlink) and purple (folHlink) from overriding run-level colors.
        # This ensures hyperlinked text preserves its original Typst color.
        self._neutralize_theme_hyperlink_colors(prs)

        # Save
        prs.save(output_path)

        if self.config.verbose:
            print(f"Saved PPTX to {output_path}")

        return output_path

    def _neutralize_theme_hyperlink_colors(self, prs):
        """Override theme-level hyperlink colors to prevent PowerPoint defaults.

        PowerPoint themes define hlink (hyperlink) and folHlink (followed hyperlink)
        colors that override run-level solidFill when a run has hlinkClick. The
        default theme uses blue (0000FF) and purple (800080), which causes
        hyperlinked text to appear blue even when we set explicit solidFill.

        This method modifies the theme XML to set both hyperlink colors to black,
        ensuring that our run-level solidFill colors take effect.
        """
        from lxml import etree

        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        for master in prs.slide_masters:
            for rel in master.part.rels.values():
                if 'theme' in str(rel.reltype):
                    theme_part = rel.target_part
                    theme_xml = etree.fromstring(theme_part.blob)

                    modified = False
                    for tag_name in ('hlink', 'folHlink'):
                        elem = theme_xml.find(f'.//{{{ns_a}}}{tag_name}')
                        if elem is not None:
                            color = elem.find(f'{{{ns_a}}}srgbClr')
                            if color is not None:
                                color.set('val', '000000')
                                modified = True

                    if modified:
                        theme_part._blob = etree.tostring(
                            theme_xml, xml_declaration=True,
                            encoding='UTF-8', standalone=True
                        )

    def _apply_links(self, page_data: PageData, slide):
        """Apply hyperlinks from SVG <a> elements to overlapping shapes.

        For each link region, find text shapes that overlap with the link's
        bounding rect and add a hyperlink to all runs within the overlap region.

        Only external links (http://, https://, mailto:) are applied as
        clickable hyperlinks. Internal document links (fragment anchors like
        #page-3 from Typst's #link(<label>)) are skipped since they don't
        map to valid PPTX hyperlink targets.

        The link styling preserves the original text color from Typst rather
        than applying the default PPTX blue+underline link style.
        """
        from pptx.util import Emu
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from lxml import etree

        for link in page_data.links:
            if not link.href:
                continue

            # Skip internal document links (fragment anchors, relative refs)
            # Only process external URLs (http, https, mailto)
            href = link.href
            if not href.startswith(('http://', 'https://', 'mailto:')):
                continue

            # Link bounding box in EMU
            link_x = link.x * EMU_PER_PX
            link_y = link.y * EMU_PER_PX
            link_r = (link.x + link.width) * EMU_PER_PX
            link_b = (link.y + link.height) * EMU_PER_PX
            link_cx = (link_x + link_r) / 2  # Center of link region

            # Find shapes that overlap with this link region
            linked = False
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                # Shape bounding box
                sx = shape.left
                sy = shape.top
                sr = sx + shape.width
                sb = sy + shape.height

                # Check overlap
                if sx < link_r and sr > link_x and sy < link_b and sb > link_y:
                    try:
                        # Add relationship once per link per slide
                        rId = slide.part.relate_to(
                            href,
                            RT.HYPERLINK,
                            is_external=True,
                        )
                        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

                        # Estimate per-run position within the textbox
                        # to find which runs overlap with the link region
                        runs_info = []
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                runs_info.append(run)

                        if not runs_info:
                            continue

                        # Approximate run positions based on character widths
                        total_chars = sum(len(r.text) for r in runs_info)
                        if total_chars == 0:
                            continue

                        shape_text_width = shape.width
                        char_width = shape_text_width / total_chars if total_chars > 0 else 0
                        current_x = sx

                        for run in runs_info:
                            run_width = len(run.text) * char_width
                            run_x = current_x
                            run_r = run_x + run_width
                            run_center = (run_x + run_r) / 2
                            current_x = run_r

                            # Check if run CENTER is inside the link region
                            if link_x <= run_center <= link_r and run.text.strip():
                                rPr = run._r.get_or_add_rPr()
                                # Build hlinkClick with proper namespace
                                hlink_xml = f'<a:hlinkClick xmlns:a="{ns_a}" xmlns:r="{ns_r}" r:id="{rId}"/>'
                                hlink_elem = etree.fromstring(hlink_xml)
                                rPr.append(hlink_elem)

                                # Explicitly suppress PowerPoint's default hyperlink
                                # styling (blue color + underline) by:
                                # 1. Setting underline to "none" to prevent auto-underline
                                rPr.set('u', 'none')
                                # 2. If no explicit color is set, set the original
                                #    text color to prevent theme hyperlink color override
                                existing_fill = rPr.find(f'{{{ns_a}}}solidFill')
                                if existing_fill is None:
                                    # Preserve original color from Typst
                                    # Default to black if no color was set
                                    fill_xml = f'<a:solidFill xmlns:a="{ns_a}"><a:srgbClr val="000000"/></a:solidFill>'
                                    rPr.insert(0, etree.fromstring(fill_xml))
                                linked = True

                    except Exception as e:
                        if self.config.verbose:
                            print(f"Warning: Failed to add hyperlink '{href}': {e}")

    def _convert_page_shapes(
        self,
        page_data: PageData,
        slide,
        page_width: float,
        page_height: float,
        svg_data: TypstSVGData,
        prefix_to_style: Dict[str, str],
    ):
        """Convert SVG shapes to PowerPoint native shapes."""
        for shape in page_data.shapes:
            try:
                self._convert_shape(shape, slide, page_data, svg_data)
            except Exception as e:
                if self.config.verbose:
                    print(f"Warning: Failed to convert shape {shape.tag}: {e}")
                # Rasterization fallback: render the failed shape as a PNG image
                try:
                    self._rasterize_shape_fallback(shape, slide, page_data)
                except Exception as e2:
                    if self.config.verbose:
                        print(f"Warning: Rasterization fallback also failed: {e2}")

    def _convert_shape(
        self,
        shape: ShapeElement,
        slide,
        page_data: PageData,
        svg_data: TypstSVGData,
    ):
        """Convert a single SVG shape to a PowerPoint shape."""
        elem = shape.element
        tag = shape.tag

        # Apply accumulated transforms to get position
        transforms = shape.transform_matrix
        # Parse transforms to get dx, dy, sx, sy
        total_dx, total_dy, total_sx, total_sy = 0, 0, 1, 1
        i = 0
        while i + 4 < len(transforms):
            dx, dy, sx, sy, rot = transforms[i:i+5]
            total_dx += dx
            total_dy += dy
            total_sx *= sx
            total_sy *= sy
            i += 5

        # Subtract page offset for y
        total_dy -= page_data.y_offset

        if tag == 'rect':
            self._add_rect(elem, slide, total_dx, total_dy, total_sx, total_sy)
        elif tag == 'circle':
            self._add_circle(elem, slide, total_dx, total_dy, total_sx, total_sy)
        elif tag == 'ellipse':
            self._add_ellipse(elem, slide, total_dx, total_dy, total_sx, total_sy)
        elif tag == 'line':
            self._add_line(elem, slide, total_dx, total_dy, total_sx, total_sy)
        elif tag == 'image':
            self._add_image(elem, slide, total_dx, total_dy, total_sx, total_sy, svg_data)
        elif tag == 'path':
            self._add_path(elem, slide, total_dx, total_dy, total_sx, total_sy, page_data)
        elif tag in ('polygon', 'polyline'):
            self._add_polygon(elem, slide, total_dx, total_dy, total_sx, total_sy)

    def _add_rect(self, elem, slide, dx, dy, sx, sy):
        """Add a rectangle shape to the slide."""
        from pptx.util import Emu
        from pptx.dml.color import RGBColor

        x = float(elem.get('x', '0')) * sx + dx
        y = float(elem.get('y', '0')) * sy + dy
        w = float(elem.get('width', '0')) * abs(sx)
        h = float(elem.get('height', '0')) * abs(sy)

        if w <= 0 or h <= 0:
            return

        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Emu(int(x * EMU_PER_PX)),
            Emu(int(y * EMU_PER_PX)),
            Emu(int(w * EMU_PER_PX)),
            Emu(int(h * EMU_PER_PX)),
        )

        # Set fill (with alpha support)
        fill_color = elem.get('fill', 'none')
        if fill_color and fill_color != 'none':
            color_alpha = _parse_color_with_alpha(fill_color)
            if color_alpha:
                hex_color, alpha = color_alpha
                # Also check fill-opacity / opacity attributes
                opacity = elem.get('fill-opacity', elem.get('opacity', None))
                if opacity:
                    try:
                        alpha = min(alpha, float(opacity))
                    except ValueError:
                        pass
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(hex_color)
                if alpha < 1.0:
                    _apply_alpha_to_shape(shape, alpha)
        else:
            shape.fill.background()

        # Set stroke
        stroke = elem.get('stroke', 'none')
        stroke_width = float(elem.get('stroke-width', '0'))
        if stroke and stroke != 'none' and stroke_width > 0:
            hex_color = _parse_color(stroke)
            if hex_color:
                shape.line.color.rgb = RGBColor.from_string(hex_color)
                shape.line.width = Emu(int(stroke_width * EMU_PER_PX * abs(sx)))
        else:
            shape.line.fill.background()

    def _add_circle(self, elem, slide, dx, dy, sx, sy):
        """Add a circle/ellipse shape to the slide."""
        from pptx.util import Emu
        from pptx.dml.color import RGBColor

        cx = float(elem.get('cx', '0')) * sx + dx
        cy = float(elem.get('cy', '0')) * sy + dy
        r = float(elem.get('r', '0'))

        x = cx - r * abs(sx)
        y = cy - r * abs(sy)
        w = 2 * r * abs(sx)
        h = 2 * r * abs(sy)

        if w <= 0 or h <= 0:
            return

        shape = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            Emu(int(x * EMU_PER_PX)),
            Emu(int(y * EMU_PER_PX)),
            Emu(int(w * EMU_PER_PX)),
            Emu(int(h * EMU_PER_PX)),
        )

        fill_color = elem.get('fill', 'none')
        if fill_color and fill_color != 'none':
            color_alpha = _parse_color_with_alpha(fill_color)
            if color_alpha:
                hex_color, alpha = color_alpha
                opacity = elem.get('fill-opacity', elem.get('opacity', None))
                if opacity:
                    try:
                        alpha = min(alpha, float(opacity))
                    except ValueError:
                        pass
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(hex_color)
                if alpha < 1.0:
                    _apply_alpha_to_shape(shape, alpha)
        else:
            shape.fill.background()

        stroke = elem.get('stroke', 'none')
        stroke_width = float(elem.get('stroke-width', '0'))
        if stroke and stroke != 'none' and stroke_width > 0:
            hex_color = _parse_color(stroke)
            if hex_color:
                shape.line.color.rgb = RGBColor.from_string(hex_color)
                shape.line.width = Emu(int(stroke_width * EMU_PER_PX))
        else:
            shape.line.fill.background()

    def _add_ellipse(self, elem, slide, dx, dy, sx, sy):
        """Add an ellipse shape to the slide."""
        from pptx.util import Emu
        from pptx.dml.color import RGBColor

        cx = float(elem.get('cx', '0')) * sx + dx
        cy = float(elem.get('cy', '0')) * sy + dy
        rx = float(elem.get('rx', '0')) * abs(sx)
        ry = float(elem.get('ry', '0')) * abs(sy)

        x = cx - rx
        y = cy - ry
        w = 2 * rx
        h = 2 * ry

        if w <= 0 or h <= 0:
            return

        shape = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            Emu(int(x * EMU_PER_PX)),
            Emu(int(y * EMU_PER_PX)),
            Emu(int(w * EMU_PER_PX)),
            Emu(int(h * EMU_PER_PX)),
        )

        fill_color = elem.get('fill', 'none')
        if fill_color and fill_color != 'none':
            color_alpha = _parse_color_with_alpha(fill_color)
            if color_alpha:
                hex_color, alpha = color_alpha
                opacity = elem.get('fill-opacity', elem.get('opacity', None))
                if opacity:
                    try:
                        alpha = min(alpha, float(opacity))
                    except ValueError:
                        pass
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(hex_color)
                if alpha < 1.0:
                    _apply_alpha_to_shape(shape, alpha)
        else:
            shape.fill.background()

    def _add_line(self, elem, slide, dx, dy, sx, sy):
        """Add a line shape to the slide."""
        from pptx.util import Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        x1 = float(elem.get('x1', '0')) * sx + dx
        y1 = float(elem.get('y1', '0')) * sy + dy
        x2 = float(elem.get('x2', '0')) * sx + dx
        y2 = float(elem.get('y2', '0')) * sy + dy

        # Line connector
        shape = slide.shapes.add_connector(
            1,  # MSO_CONNECTOR.STRAIGHT
            Emu(int(x1 * EMU_PER_PX)),
            Emu(int(y1 * EMU_PER_PX)),
            Emu(int(x2 * EMU_PER_PX)),
            Emu(int(y2 * EMU_PER_PX)),
        )

        stroke = elem.get('stroke', '#000000')
        if stroke and stroke != 'none':
            hex_color = _parse_color(stroke)
            if hex_color:
                shape.line.color.rgb = RGBColor.from_string(hex_color)

        stroke_width = float(elem.get('stroke-width', '1'))
        shape.line.width = Emu(int(stroke_width * EMU_PER_PX * abs(sx)))

    def _rasterize_image_to_png(self, image_data: bytes, image_format: str,
                                width_px: int = 0, height_px: int = 0) -> bytes:
        """Rasterize SVG or PDF image data to PNG using the typst Python package.

        Uses typst.compile(format="png") with transparent background (fill: none)
        to convert SVG/PDF images to high-quality PNG for embedding in PPTX.

        Args:
            image_data: Raw image bytes (SVG or PDF)
            image_format: Image format string (e.g. 'svg+xml', 'svg', 'pdf')
            width_px: Target width in pixels (0 = use default 400px)
            height_px: Target height in pixels (0 = auto from width)

        Returns:
            PNG image data as bytes, or empty bytes on failure
        """
        try:
            import typst as typst_py

            # Map format to file extension
            ext_map = {'svg': '.svg', 'svg+xml': '.svg', 'pdf': '.pdf'}
            ext = ext_map.get(image_format, f'.{image_format}')
            img_filename = f"image{ext}"

            # Default width if not specified
            if width_px <= 0:
                width_px = 400
            # Use auto height if not specified (typst will compute aspect ratio)
            height_spec = f", height: {height_px}pt" if height_px > 0 else ""

            with tempfile.TemporaryDirectory() as td:
                img_path = os.path.join(td, img_filename)
                typ_path = os.path.join(td, "main.typ")

                with open(img_path, 'wb') as f:
                    f.write(image_data)

                # Use fill: none for transparent background
                typ_source = (
                    f'#set page(width: {width_px}pt{height_spec}, '
                    f'margin: 0pt, fill: none)\n'
                    f'#image("{img_filename}", width: 100%'
                    f'{", height: 100%" if height_px > 0 else ""}'
                    f', fit: "contain")'
                )
                with open(typ_path, 'w') as f:
                    f.write(typ_source)

                dpi = self.config.raster_dpi if hasattr(self.config, 'raster_dpi') else 300
                result = typst_py.compile(typ_path, format="png", ppi=float(dpi))
                if isinstance(result, list):
                    result = result[0]
                return result

        except Exception as e:
            if self.config.verbose:
                print(f"Warning: Image rasterization failed ({image_format}): {e}")
            return b''

    def _rasterize_shape_fallback(self, shape, slide, page_data):
        """Rasterize a failed SVG shape to PNG and add it as an image.

        When native shape conversion fails, this method serializes the SVG element
        back to SVG markup and rasterizes it to PNG via the typst Python package.
        """
        from lxml import etree
        from pptx.util import Emu
        from io import BytesIO

        elem = shape.element
        transforms = shape.transform_matrix

        # Get bounding box from element attributes or transform
        total_dx, total_dy, total_sx, total_sy = 0, 0, 1, 1
        i = 0
        while i + 4 < len(transforms):
            dx, dy, sx, sy, rot = transforms[i:i+5]
            total_dx += dx
            total_dy += dy
            total_sx *= sx
            total_sy *= sy
            i += 5
        total_dy -= page_data.y_offset

        # Try to determine element dimensions
        w = float(elem.get('width', 0) or 0) * abs(total_sx)
        h = float(elem.get('height', 0) or 0) * abs(total_sy)
        if w <= 0 or h <= 0:
            # Try bounding box from path/shape attributes
            x1 = float(elem.get('x1', elem.get('cx', elem.get('x', 0))))
            y1 = float(elem.get('y1', elem.get('cy', elem.get('y', 0))))
            x2 = float(elem.get('x2', elem.get('r', 50)))
            y2 = float(elem.get('y2', elem.get('r', 50)))
            w = max(abs(x2 - x1), abs(x2)) * abs(total_sx) if w <= 0 else w
            h = max(abs(y2 - y1), abs(y2)) * abs(total_sy) if h <= 0 else h
        if w <= 0:
            w = 100
        if h <= 0:
            h = 100

        # Serialize element to SVG string
        elem_svg = etree.tostring(elem, encoding='unicode')
        svg_wrapper = (
            f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'width="{w}" height="{h}" viewBox="0 0 {w} {h}">'
            f'{elem_svg}</svg>'
        )

        png_data = self._rasterize_image_to_png(
            svg_wrapper.encode('utf-8'), 'svg',
            width_px=int(w * 2), height_px=int(h * 2)
        )
        if not png_data:
            return

        # Convert SVG coordinates to EMU
        slide_width = slide.slide_layout.slide_master.slide_width if hasattr(slide, 'slide_layout') else 9144000
        slide_height = slide.slide_layout.slide_master.slide_height if hasattr(slide, 'slide_layout') else 5143500
        page_width = getattr(page_data, 'width', 842)
        page_height = getattr(page_data, 'height', 474)
        scale_x = slide_width / page_width
        scale_y = slide_height / page_height

        left_emu = int(total_dx * scale_x)
        top_emu = int(total_dy * scale_y)
        width_emu = int(w * scale_x)
        height_emu = int(h * scale_y)

        img_stream = BytesIO(png_data)
        slide.shapes.add_picture(img_stream, left_emu, top_emu, width_emu, height_emu)

    def _add_image(self, elem, slide, dx, dy, sx, sy, svg_data):
        """Add an image to the slide.

        Handles PNG, JPEG, GIF images directly.
        SVG and PDF images are rasterized to PNG via the typst Python package
        with transparent background before embedding.
        """
        from pptx.util import Emu
        import base64

        href = elem.get(f'{{{XLINK_NS}}}href') or elem.get('href', '')
        if not href:
            return

        x = float(elem.get('x', '0')) * sx + dx
        y = float(elem.get('y', '0')) * sy + dy
        w = float(elem.get('width', '0')) * abs(sx)
        h = float(elem.get('height', '0')) * abs(sy)

        if w <= 0 or h <= 0:
            return

        left_emu = Emu(int(x * EMU_PER_PX))
        top_emu = Emu(int(y * EMU_PER_PX))
        width_emu = Emu(int(w * EMU_PER_PX))
        height_emu = Emu(int(h * EMU_PER_PX))

        if href.startswith('data:'):
            # Data URI
            match = re.match(r'data:image/([^;]+);base64,(.*)', href, re.DOTALL)
            if match:
                fmt = match.group(1).lower()
                data = base64.b64decode(match.group(2))

                if fmt in ('svg+xml', 'svg', 'pdf'):
                    # SVG/PDF images: rasterize to PNG via typst
                    png_data = self._rasterize_image_to_png(
                        data, fmt,
                        width_px=int(w * 2), height_px=int(h * 2)
                    )
                    if not png_data:
                        return
                    img_stream = BytesIO(png_data)
                else:
                    # PNG, JPEG, GIF, etc. - pass directly
                    img_stream = BytesIO(data)

                slide.shapes.add_picture(
                    img_stream, left_emu, top_emu, width_emu, height_emu
                )
        else:
            # External file
            img_path = Path(href)
            if not img_path.is_absolute():
                # Try relative to SVG directory
                img_path = Path(svg_data.defs.get('__svg_dir', '')) / href
            if img_path.exists():
                suffix = img_path.suffix.lower()
                if suffix in ('.svg', '.pdf'):
                    # Rasterize SVG/PDF file to PNG via typst
                    fmt = 'svg' if suffix == '.svg' else 'pdf'
                    with open(img_path, 'rb') as f:
                        file_bytes = f.read()
                    png_data = self._rasterize_image_to_png(
                        file_bytes, fmt,
                        width_px=int(w * 2), height_px=int(h * 2)
                    )
                    if not png_data:
                        return
                    img_stream = BytesIO(png_data)
                    slide.shapes.add_picture(
                        img_stream, left_emu, top_emu, width_emu, height_emu
                    )
                else:
                    slide.shapes.add_picture(
                        str(img_path), left_emu, top_emu, width_emu, height_emu
                    )

    def _add_path(self, elem, slide, dx, dy, sx, sy, page_data):
        """Add a path shape to the slide as native DrawingML custom geometry.

        Uses the ppt-master path pipeline:
        1. Parse SVG path d attribute into commands
        2. Convert to absolute coordinates
        3. Normalize all curves to cubic beziers
        4. Generate DrawingML <a:custGeom> XML
        5. Inject raw XML into the slide via lxml
        """
        from lxml import etree

        path_d = elem.get('d', '')
        if not path_d:
            return

        # Use ppt-master's path conversion pipeline
        try:
            from ..scripts.svg_to_shapes import (
                parse_svg_path, svg_path_to_absolute,
                normalize_path_commands, path_commands_to_drawingml,
            )
        except ImportError:
            if self.config.verbose:
                print("Warning: svg_to_shapes not available, skipping path")
            return

        # Parse and convert path
        commands = parse_svg_path(path_d)
        commands = svg_path_to_absolute(commands)
        commands = normalize_path_commands(commands)

        if not commands:
            return

        # Handle transform on the path element itself
        path_transform = elem.get('transform', '')
        extra_dx, extra_dy = 0.0, 0.0
        if path_transform:
            import re as _re
            t_match = _re.search(r'translate\(\s*([-\d.e+]+)[\s,]+([-\d.e+]+)\s*\)', path_transform)
            if t_match:
                extra_dx = float(t_match.group(1))
                extra_dy = float(t_match.group(2))

        # Convert to DrawingML path commands
        path_xml, min_x, min_y, width, height = path_commands_to_drawingml(
            commands,
            offset_x=dx + extra_dx * abs(sx),
            offset_y=dy + extra_dy * abs(sy),
            scale_x=sx,
            scale_y=sy,
        )

        if not path_xml or width <= 0 or height <= 0:
            return

        # Build fill XML
        fill_color = elem.get('fill', 'none')
        fill_xml = ''
        if fill_color and fill_color != 'none':
            if fill_color.startswith('url(#'):
                # Gradient fill reference
                fill_xml = self._build_gradient_fill_xml(fill_color, page_data)
                if not fill_xml:
                    fill_xml = '<a:noFill/>'
            else:
                color_alpha = _parse_color_with_alpha(fill_color)
                if color_alpha:
                    hex_color, color_a = color_alpha
                    # Also check fill-opacity / opacity attributes
                    opacity = elem.get('fill-opacity', elem.get('opacity', '1'))
                    try:
                        opacity_val = float(opacity)
                    except ValueError:
                        opacity_val = 1.0
                    # Combine color alpha with element opacity
                    combined_alpha = color_a * opacity_val
                    alpha = int(combined_alpha * 100000)
                    if alpha < 100000:
                        fill_xml = f'<a:solidFill><a:srgbClr val="{hex_color}"><a:alpha val="{alpha}"/></a:srgbClr></a:solidFill>'
                    else:
                        fill_xml = f'<a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
                else:
                    fill_xml = '<a:noFill/>'
        else:
            fill_xml = '<a:noFill/>'

        # Build stroke XML (with alpha support)
        stroke_color = elem.get('stroke', 'none')
        stroke_width = float(elem.get('stroke-width', '0') or '0')
        stroke_xml = ''
        if stroke_color and stroke_color != 'none' and stroke_width > 0:
            stroke_alpha = _parse_color_with_alpha(stroke_color)
            if stroke_alpha:
                hex_stroke, s_alpha = stroke_alpha
                # Also check stroke-opacity
                s_opacity = elem.get('stroke-opacity', elem.get('opacity', '1'))
                try:
                    s_alpha = s_alpha * float(s_opacity)
                except ValueError:
                    pass
                sw_emu = int(stroke_width * abs(sx) * EMU_PER_PX)
                if s_alpha < 1.0:
                    s_alpha_val = int(s_alpha * 100000)
                    stroke_xml = f'<a:ln w="{sw_emu}"><a:solidFill><a:srgbClr val="{hex_stroke}"><a:alpha val="{s_alpha_val}"/></a:srgbClr></a:solidFill></a:ln>'
                else:
                    stroke_xml = f'<a:ln w="{sw_emu}"><a:solidFill><a:srgbClr val="{hex_stroke}"/></a:solidFill></a:ln>'

        # Build geometry XML
        w_emu = int(width * EMU_PER_PX)
        h_emu = int(height * EMU_PER_PX)

        geom_xml = f'''<a:custGeom>
<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
<a:rect l="l" t="t" r="r" b="b"/>
<a:pathLst><a:path w="{w_emu}" h="{h_emu}">
{path_xml}
</a:path></a:pathLst>
</a:custGeom>'''

        # Build the full <p:sp> shape XML
        shape_id = len(slide.shapes) + 100  # Avoid ID conflicts
        off_x_emu = int(min_x * EMU_PER_PX)
        off_y_emu = int(min_y * EMU_PER_PX)

        # Ensure non-negative position
        if off_x_emu < 0:
            off_x_emu = 0
        if off_y_emu < 0:
            off_y_emu = 0

        shape_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
<p:nvSpPr>
<p:cNvPr id="{shape_id}" name="Freeform {shape_id}"/>
<p:cNvSpPr/><p:nvPr/>
</p:nvSpPr>
<p:spPr>
<a:xfrm><a:off x="{off_x_emu}" y="{off_y_emu}"/><a:ext cx="{w_emu}" cy="{h_emu}"/></a:xfrm>
{geom_xml}
{fill_xml}
{stroke_xml}
</p:spPr>
</p:sp>'''

        # Parse and inject into slide's spTree
        try:
            shape_elem = etree.fromstring(shape_xml)
            sp_tree = slide.shapes._spTree
            sp_tree.append(shape_elem)
        except Exception as e:
            if self.config.verbose:
                print(f"Warning: Failed to inject path shape XML: {e}")

    def _build_gradient_fill_xml(self, fill_ref: str, page_data) -> Optional[str]:
        """Build DrawingML gradient fill XML from SVG url(#id) reference.

        Resolves the gradient from SVG defs, extracts stop colors,
        and generates <a:gradFill> XML for DrawingML.

        Args:
            fill_ref: Fill attribute value like 'url(#gradientId)'
            page_data: PageData for accessing SVG defs

        Returns:
            DrawingML gradient fill XML string, or None if not resolvable
        """
        import re as _re

        # Extract gradient ID from url(#id)
        m = _re.match(r'url\(#(.+?)\)', fill_ref)
        if not m:
            return None
        grad_id = m.group(1)

        # Look up in SVG defs
        svg_data = self._svg_data
        if not svg_data or not svg_data.defs:
            return None

        grad_elem = svg_data.defs.get(grad_id)
        if grad_elem is None:
            return None

        # Follow href to find base gradient with stops
        ns_xlink = '{http://www.w3.org/1999/xlink}'
        href = grad_elem.get('href') or grad_elem.get(f'{ns_xlink}href')
        stops_elem = grad_elem
        if href and href.startswith('#'):
            base_id = href[1:]
            base_elem = svg_data.defs.get(base_id)
            if base_elem is not None:
                stops_elem = base_elem

        # Extract stops
        ns_svg = '{http://www.w3.org/2000/svg}'
        stops = []
        for child in stops_elem:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag == 'stop':
                offset_str = child.get('offset', '0')
                if offset_str.endswith('%'):
                    offset = float(offset_str[:-1]) / 100.0
                else:
                    offset = float(offset_str)
                color = child.get('stop-color', '#000000')
                opacity = child.get('stop-opacity', '1')
                hex_color = _parse_color(color)
                if hex_color:
                    stops.append((offset, hex_color, float(opacity)))

        if len(stops) < 2:
            return None

        # Determine gradient direction from the base or derived gradient
        x1 = float(grad_elem.get('x1', stops_elem.get('x1', '0')))
        y1 = float(grad_elem.get('y1', stops_elem.get('y1', '0')))
        x2 = float(grad_elem.get('x2', stops_elem.get('x2', '1')))
        y2 = float(grad_elem.get('y2', stops_elem.get('y2', '0')))

        # Compute angle from direction vector
        import math
        dx_grad = x2 - x1
        dy_grad = y2 - y1
        if dx_grad == 0 and dy_grad == 0:
            angle_deg = 0
        else:
            # SVG angles: 0=right, 90=down
            # DrawingML angles: in 1/60000 degrees, 0=right, positive=counterclockwise
            angle_rad = math.atan2(dy_grad, dx_grad)
            angle_deg = math.degrees(angle_rad)

        # Convert to DrawingML angle (60000ths of a degree)
        # DrawingML: 0 = left-to-right, angles are counterclockwise
        drawingml_angle = int(angle_deg * 60000)

        # For PowerPoint, limit to key gradient stops
        # Use first, last, and optionally a mid stop
        if len(stops) <= 3:
            key_stops = stops
        else:
            # Simplify: take first, mid, and last
            mid_idx = len(stops) // 2
            key_stops = [stops[0], stops[mid_idx], stops[-1]]

        # Build gsLst (gradient stop list)
        gs_list = []
        for offset, hex_color, opacity in key_stops:
            pos = int(offset * 100000)  # Position in 1/100000ths
            if opacity < 1.0:
                alpha = int(opacity * 100000)
                gs_list.append(
                    f'<a:gs pos="{pos}"><a:srgbClr val="{hex_color}">'
                    f'<a:alpha val="{alpha}"/></a:srgbClr></a:gs>'
                )
            else:
                gs_list.append(
                    f'<a:gs pos="{pos}"><a:srgbClr val="{hex_color}"/></a:gs>'
                )

        gs_xml = '\n'.join(gs_list)

        return (
            f'<a:gradFill>'
            f'<a:gsLst>{gs_xml}</a:gsLst>'
            f'<a:lin ang="{drawingml_angle}" scaled="1"/>'
            f'</a:gradFill>'
        )

    def _add_polygon(self, elem, slide, dx, dy, sx, sy):
        """Add a polygon/polyline shape."""
        # Parse points
        points_str = elem.get('points', '')
        if not points_str:
            return

        # Parse point pairs
        points = []
        for pair in re.findall(r'([-\d.e+]+)[,\s]+([-\d.e+]+)', points_str):
            px = float(pair[0]) * sx + dx
            py = float(pair[1]) * sy + dy
            points.append((px, py))

        if len(points) < 2:
            return

        # Use freeform builder from python-pptx
        from pptx.util import Emu

        # Calculate bounding box
        xs = [p[0] for p in points]
        ys = [p[1] for p in points]
        min_x, max_x = min(xs), max(xs)
        min_y, max_y = min(ys), max(ys)
        width = max_x - min_x
        height = max_y - min_y

        if width <= 0 or height <= 0:
            return

        # Build freeform shape
        freeform = slide.shapes.build_freeform(
            Emu(int(min_x * EMU_PER_PX)),
            Emu(int(min_y * EMU_PER_PX)),
        )

        # Move to first point (relative to bounding box origin)
        freeform.move_to(
            Emu(int((points[0][0] - min_x) * EMU_PER_PX)),
            Emu(int((points[0][1] - min_y) * EMU_PER_PX)),
        )

        # Line to remaining points
        for px, py in points[1:]:
            freeform.line_to(
                Emu(int((px - min_x) * EMU_PER_PX)),
                Emu(int((py - min_y) * EMU_PER_PX)),
            )

        # Close if polygon
        tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
        if tag == 'polygon':
            freeform.line_to(
                Emu(int((points[0][0] - min_x) * EMU_PER_PX)),
                Emu(int((points[0][1] - min_y) * EMU_PER_PX)),
            )

        shape = freeform.convert_to_shape(
            Emu(int(width * EMU_PER_PX)),
            Emu(int(height * EMU_PER_PX)),
        )

    def _convert_page_texts(
        self,
        page_data: PageData,
        slide,
        page_width: float,
        page_height: float,
        prefix_to_style: Dict[str, str],
    ):
        """Convert text segments to PowerPoint text boxes.

        Groups text segments on the same line into a single textbox
        with multiple runs, so that text like "Hello **world** today"
        becomes one textbox with three runs (regular, bold, regular)
        instead of three separate textboxes.

        Math formula segments are handled separately: they are grouped
        spatially into formula regions with a single textbox each.

        Inline math (math segments on the same line as text) is merged
        into the text textbox as Cambria Math runs.
        """
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_AUTO_SIZE

        # Filter out empty segments
        segments = [seg for seg in page_data.text_segments if seg.text.strip()]
        if not segments:
            return

        # Separate math segments from non-math segments
        math_segments = [s for s in segments if s.font_variant == 'math']
        text_segments = [s for s in segments if s.font_variant != 'math']

        # Classify math as inline or display using cluster-first approach
        inline_math = []
        display_math = []

        if math_segments and text_segments:
            def get_baseline(seg):
                return seg.y + seg.font_size * 0.8

            # First, cluster ALL math segments into spatial groups
            math_clusters = self._cluster_math_segments(math_segments)

            # Get text line baselines (group text segments by baseline first)
            text_baselines = []
            text_x_ranges = []
            text_by_bl = {}
            for ts in text_segments:
                ts_bl = round(get_baseline(ts), 1)
                text_by_bl.setdefault(ts_bl, []).append(ts)

            # Merge nearby text baselines
            sorted_bls = sorted(text_by_bl.keys())
            merged_text_lines = []
            for bl in sorted_bls:
                if merged_text_lines and abs(bl - merged_text_lines[-1][0]) < 3.0:
                    merged_text_lines[-1][1].extend(text_by_bl[bl])
                    merged_text_lines[-1] = (
                        sum(get_baseline(s) for s in merged_text_lines[-1][1]) / len(merged_text_lines[-1][1]),
                        merged_text_lines[-1][1]
                    )
                else:
                    merged_text_lines.append((bl, text_by_bl[bl]))

            # For each math cluster, check if it's inline with a text line
            for cluster in math_clusters:
                # Get the cluster's dominant-size baseline
                font_sizes = [s.font_size for s in cluster]
                dominant_size = max(set(round(s, 1) for s in font_sizes),
                                   key=lambda sz: sum(1 for s in font_sizes if round(s, 1) == sz))

                anchor_segs = [s for s in cluster
                               if abs(s.font_size - dominant_size) < dominant_size * 0.2]
                if not anchor_segs:
                    anchor_segs = cluster

                cluster_bl = sum(get_baseline(s) for s in anchor_segs) / len(anchor_segs)
                cluster_x_min = min(s.x for s in cluster)
                cluster_x_max = max(s.x + s.width for s in cluster)

                is_inline = False
                for text_bl, text_segs in merged_text_lines:
                    # Check baseline proximity (use dominant size tolerance)
                    if abs(cluster_bl - text_bl) < dominant_size * 0.6:
                        # Also check x-proximity: cluster should be near text
                        text_x_min = min(s.x for s in text_segs)
                        text_x_max = max(s.x + s.width for s in text_segs)
                        # Inline if cluster is within reasonable distance of text line
                        x_gap = max(0, max(cluster_x_min - text_x_max,
                                          text_x_min - cluster_x_max))
                        if x_gap < dominant_size * 5:
                            is_inline = True
                            break

                if is_inline:
                    inline_math.extend(cluster)
                else:
                    display_math.extend(cluster)
        else:
            display_math = math_segments

        # Handle inline math based on config
        if self.config.inline_math_mode == "glyph":
            # Render all inline math as glyph curves
            has_inline_glyph = any(s.glyph_uses for s in inline_math) if inline_math else False
            if has_inline_glyph:
                display_math.extend(inline_math)
                combined_segments = text_segments
            else:
                combined_segments = text_segments + inline_math
        elif self.config.inline_math_mode == "auto" and inline_math:
            # Heuristic: classify each inline math cluster as simple or complex
            # Re-cluster inline math to process per-formula
            inline_clusters = self._cluster_math_segments(inline_math) if inline_math else []
            simple_math = []
            complex_math = []
            for cluster in inline_clusters:
                if _is_simple_inline_math(cluster):
                    simple_math.extend(cluster)
                else:
                    complex_math.extend(cluster)
            # Simple math rendered as text, complex as glyph curves
            combined_segments = text_segments + simple_math
            if complex_math:
                has_glyph = any(s.glyph_uses for s in complex_math)
                if has_glyph:
                    display_math.extend(complex_math)
                else:
                    combined_segments.extend(complex_math)
        else:
            # "text" mode: merge inline math with text as Cambria Math runs
            combined_segments = text_segments + inline_math

        # Process combined text + inline math segments (line-based grouping)
        self._render_text_groups(combined_segments, slide, page_width=page_width)

        # Process display math segments
        if display_math:
            if self.config.display_math_mode == "glyph":
                # Render as native curves (glyph outlines)
                has_glyph_data = any(s.glyph_uses for s in display_math)
                if has_glyph_data:
                    self._render_math_as_curves(display_math, slide)
                else:
                    # Fallback to text-based rendering if no glyph data available
                    self._render_math_groups(display_math, text_segments, slide)
            elif self.config.display_math_mode == "auto":
                # Heuristic: classify each display math cluster as simple or complex
                display_clusters = self._cluster_math_segments(display_math)
                simple_display = []
                complex_display = []
                for cluster in display_clusters:
                    if _is_simple_inline_math(cluster):
                        simple_display.extend(cluster)
                    else:
                        complex_display.extend(cluster)
                # Render simple display math as text
                if simple_display:
                    self._render_math_groups(simple_display, text_segments, slide)
                # Render complex display math as glyph curves
                if complex_display:
                    has_glyph_data = any(s.glyph_uses for s in complex_display)
                    if has_glyph_data:
                        self._render_math_as_curves(complex_display, slide)
                    else:
                        self._render_math_groups(complex_display, text_segments, slide)
            else:
                # "text" mode: Render display math as text (Cambria Math runs)
                self._render_math_groups(display_math, text_segments, slide)

    def _detect_alignment(self, line_segs: list, page_width: float) -> str:
        """Detect text alignment from segment positions relative to page width.

        Returns 'left', 'center', or 'right'.
        """
        if not line_segs or page_width <= 0:
            return 'left'

        line_left = min(s.x for s in line_segs)
        line_right = max(s.x + s.width for s in line_segs)
        line_center = (line_left + line_right) / 2
        line_width = line_right - line_left
        page_center = page_width / 2

        # Don't try to detect alignment for lines that fill most of the width
        # (these are body text, alignment is ambiguous)
        if line_width > page_width * 0.7:
            return 'left'

        # Check center alignment: line center is near page center
        # and there's significant margin on both sides
        center_offset = abs(line_center - page_center)
        if center_offset < page_width * 0.03:  # within 3% of center
            left_margin = line_left
            right_margin = page_width - line_right
            # Both margins should be substantial and roughly equal
            if left_margin > page_width * 0.1 and right_margin > page_width * 0.1:
                margin_ratio = min(left_margin, right_margin) / max(left_margin, right_margin)
                if margin_ratio > 0.6:  # margins within 40% of each other
                    return 'center'

        # Check right alignment: right edge near page right margin
        # and substantial left margin
        right_margin = page_width - line_right
        if right_margin < page_width * 0.08 and line_left > page_width * 0.25:
            return 'right'

        return 'left'

    def _render_text_groups(self, segments: list, slide, page_width: float = 0):
        """Render regular (non-math) text segments grouped by line.

        Detects multi-line paragraphs (like lorem text) and merges them into
        single textboxes with word wrapping enabled.
        """
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_AUTO_SIZE

        if not segments:
            return

        # Consolidate dot leader characters into single segments
        # TOC entries have many individual '.' segments that should be
        # collapsed into one segment for cleaner rendering
        segments = _consolidate_dot_leaders(segments)

        line_groups = _group_segments_by_line(segments)

        # Merge inline math sub/superscripts into their adjacent text lines
        line_groups = _merge_inline_math_subscripts(line_groups)

        # Detect and merge paragraph lines
        merged_groups = _merge_paragraph_lines(line_groups, page_width=page_width)

        for group in merged_groups:
            if group['type'] == 'paragraph':
                self._add_paragraph_textbox(slide, group, page_width=page_width)
            else:
                line_segs = group['segments']
                # Sort segments left-to-right within the line
                line_segs.sort(key=lambda s: s.x)

                # Compute the merged bounding box for all segments in this line
                merged_x = min(s.x for s in line_segs)
                merged_y = min(s.y for s in line_segs)
                merged_right = max(s.x + s.width for s in line_segs)
                merged_bottom = max(s.y + s.height for s in line_segs)
                merged_w = merged_right - merged_x
                merged_h = merged_bottom - merged_y

                # Ensure minimum dimensions
                if merged_w < 1:
                    total_text_len = sum(len(s.text) for s in line_segs)
                    avg_font_size = sum(s.font_size for s in line_segs) / len(line_segs)
                    merged_w = total_text_len * avg_font_size * 0.6
                if merged_h < 1:
                    max_font_size = max(s.font_size for s in line_segs)
                    merged_h = max_font_size * 1.3

                # Detect alignment
                alignment = self._detect_alignment(line_segs, page_width)

                self._add_textbox(slide, line_segs, merged_x, merged_y, merged_w, merged_h,
                                  alignment=alignment)

    def _add_textbox(self, slide, line_segs: list, x: float, y: float, w: float, h: float,
                     alignment: str = 'left'):
        """Add a textbox with multiple runs to the slide."""
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

        # Convert to EMU
        left = Emu(int(x * EMU_PER_PX))
        top = Emu(int(y * EMU_PER_PX))
        width = Emu(int(w * EMU_PER_PX))
        height = Emu(int(h * EMU_PER_PX))

        # Ensure positive values
        if left < 0:
            left = Emu(0)
        if top < 0:
            top = Emu(0)
        if width <= 0:
            width = Emu(914400)  # 1 inch default
        if height <= 0:
            height = Emu(457200)  # 0.5 inch default

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = False

        # Precompute baseline info for superscript/subscript detection
        # Find the dominant (most common / largest) font size for anchor baseline
        def get_baseline(seg):
            return seg.y + seg.font_size * 0.8

        non_math_segs = [s for s in line_segs if s.font_variant != 'math']
        if non_math_segs:
            dominant_size = max(s.font_size for s in non_math_segs)
            anchor_segs = [s for s in non_math_segs if s.font_size >= dominant_size * 0.9]
        else:
            dominant_size = max(s.font_size for s in line_segs) if line_segs else 0
            anchor_segs = [s for s in line_segs if s.font_size >= dominant_size * 0.9]

        if anchor_segs:
            anchor_bl = sum(get_baseline(s) for s in anchor_segs) / len(anchor_segs)
        else:
            anchor_bl = sum(get_baseline(s) for s in line_segs) / len(line_segs) if line_segs else 0

        # Create one paragraph with multiple runs (one per segment)
        p = tf.paragraphs[0]

        # Set paragraph alignment
        if alignment == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif alignment == 'right':
            p.alignment = PP_ALIGN.RIGHT

        for i, seg in enumerate(line_segs):
            run = p.add_run()

            # Detect gaps between consecutive segments and insert spaces
            text = seg.text
            if i > 0:
                prev = line_segs[i - 1]
                prev_text = prev.text
                gap = seg.x - (prev.x + prev.width)
                avg_font = (seg.font_size + prev.font_size) / 2
                if gap > avg_font * 0.15:
                    if not prev_text.endswith(' ') and not text.startswith(' '):
                        text = " " + text

            run.text = text

            # Font properties
            font_props = _font_variant_to_props(seg.font_variant)
            run.font.bold = font_props['bold']
            run.font.italic = font_props['italic']
            run.font.name = font_props['font_family']

            # Font size
            font_size_pt = seg.font_size * 0.75  # px to pt
            if font_size_pt > 0:
                run.font.size = Pt(font_size_pt)

            # Detect superscript/subscript for inline math segments
            # PPTX baseline convention: positive = raised (superscript),
            # negative = lowered (subscript).
            # PPTX baseline already visually shrinks sup/sub text,
            # so we use the dominant (parent) font size, not the segment's own size.
            if dominant_size > 0 and seg.font_variant == 'math':
                size_ratio = seg.font_size / dominant_size
                if size_ratio < 0.85:
                    seg_bl = get_baseline(seg)
                    bl_diff = seg_bl - anchor_bl

                    if bl_diff < -dominant_size * 0.08:
                        # Superscript (baseline above anchor → raised text)
                        offset_pct = max(20, min(50, int(-bl_diff / dominant_size * 100)))
                        self._set_run_baseline(run, offset_pct)  # positive = raised
                        run.font.size = Pt(dominant_size * 0.75)
                    elif bl_diff > dominant_size * 0.08:
                        # Subscript (baseline below anchor → lowered text)
                        offset_pct = max(20, min(50, int(bl_diff / dominant_size * 100)))
                        self._set_run_baseline(run, -offset_pct)  # negative = lowered
                        run.font.size = Pt(dominant_size * 0.75)

            # Color
            hex_color = _parse_color(seg.fill_color)
            if hex_color:
                run.font.color.rgb = RGBColor.from_string(hex_color)

        # Set text frame auto-size
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # Remove margins
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

    def _add_paragraph_textbox(self, slide, para_group: dict, page_width: float = 0):
        """Add a multi-line paragraph textbox with word wrapping.

        This handles cases like #lorem(200) where the text wraps across
        multiple lines. Instead of creating one textbox per line, we create
        a single textbox with all lines joined together and word wrapping enabled.
        """
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

        lines = para_group['lines']
        all_segs = para_group['segments']
        line_spacing = para_group.get('line_spacing', 0)

        if not all_segs:
            return

        # Compute bounding box across all lines
        bb_x = min(s.x for s in all_segs)
        bb_y = min(s.y for s in all_segs)
        bb_right = max(s.x + s.width for s in all_segs)
        bb_bottom = max(s.y + s.height for s in all_segs)
        bb_w = bb_right - bb_x
        bb_h = bb_bottom - bb_y

        # Ensure minimum dimensions
        if bb_w < 1:
            avg_font = sum(s.font_size for s in all_segs) / len(all_segs)
            bb_w = max(len(s.text) for s in all_segs) * avg_font * 0.6
        if bb_h < 1:
            max_font = max(s.font_size for s in all_segs)
            bb_h = max_font * len(lines) * 1.5

        # Convert to EMU
        left = Emu(max(0, int(bb_x * EMU_PER_PX)))
        top = Emu(max(0, int(bb_y * EMU_PER_PX)))
        width = Emu(max(int(bb_w * EMU_PER_PX), 9525))
        height = Emu(max(int(bb_h * EMU_PER_PX), 9525))

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True  # Enable word wrapping for paragraphs

        # Helper for baseline calculation
        def get_baseline(seg):
            return seg.y + seg.font_size * 0.8

        # Detect justify alignment for paragraphs:
        # If most lines (except the last) reach near the right margin,
        # the paragraph is likely justified.
        para_alignment = None
        if page_width > 0 and len(lines) >= 2:
            full_width_lines = 0
            for li, line_segs in enumerate(lines[:-1]):  # Exclude last line
                line_right = max(s.x + s.width for s in line_segs)
                if line_right > page_width * 0.85:
                    full_width_lines += 1
            if full_width_lines >= len(lines) - 1:
                # Check if left edges also align (paragraph fills content area)
                left_edges = [min(s.x for s in line) for line in lines]
                if max(left_edges) - min(left_edges) < 10:
                    para_alignment = PP_ALIGN.JUSTIFY

        # Join all lines into a single paragraph text
        # Each line's segments are sorted left-to-right and concatenated
        p = tf.paragraphs[0]
        if para_alignment is not None:
            p.alignment = para_alignment

        for li, line_segs in enumerate(lines):
            line_segs_sorted = sorted(line_segs, key=lambda s: s.x)

            # Compute per-line anchor baseline for superscript/subscript detection
            non_math_in_line = [s for s in line_segs_sorted if s.font_variant != 'math']
            if non_math_in_line:
                line_dominant_size = max(s.font_size for s in non_math_in_line)
                line_anchors = [s for s in non_math_in_line if s.font_size >= line_dominant_size * 0.9]
            else:
                line_dominant_size = max(s.font_size for s in line_segs_sorted) if line_segs_sorted else 0
                line_anchors = [s for s in line_segs_sorted if s.font_size >= line_dominant_size * 0.9]

            if line_anchors:
                line_anchor_bl = sum(get_baseline(s) for s in line_anchors) / len(line_anchors)
            else:
                line_anchor_bl = sum(get_baseline(s) for s in line_segs_sorted) / len(line_segs_sorted) if line_segs_sorted else 0

            for si, seg in enumerate(line_segs_sorted):
                run = p.add_run()
                text = seg.text

                # Add space between segments on the same line
                if si > 0:
                    prev = line_segs_sorted[si - 1]
                    gap = seg.x - (prev.x + prev.width)
                    avg_font = (seg.font_size + prev.font_size) / 2
                    if gap > avg_font * 0.15:
                        if not prev.text.endswith(' ') and not text.startswith(' '):
                            text = " " + text

                # Add space between lines (join with space)
                if si == 0 and li > 0:
                    # Check if previous line's last segment ends with space
                    prev_line = lines[li - 1]
                    prev_last = sorted(prev_line, key=lambda s: s.x)[-1]
                    if not prev_last.text.endswith(' ') and not text.startswith(' '):
                        text = " " + text

                run.text = text

                # Font properties
                font_props = _font_variant_to_props(seg.font_variant)
                run.font.bold = font_props['bold']
                run.font.italic = font_props['italic']
                run.font.name = font_props['font_family']

                # Font size
                font_size_pt = seg.font_size * 0.75
                if font_size_pt > 0:
                    run.font.size = Pt(font_size_pt)

                # Detect superscript/subscript for inline math segments
                # PPTX baseline: positive = raised (super), negative = lowered (sub).
                if line_dominant_size > 0 and seg.font_variant == 'math':
                    size_ratio = seg.font_size / line_dominant_size
                    if size_ratio < 0.85:
                        seg_bl = get_baseline(seg)
                        bl_diff = seg_bl - line_anchor_bl

                        if bl_diff < -line_dominant_size * 0.08:
                            # Superscript (above anchor → positive baseline)
                            offset_pct = max(20, min(50, int(-bl_diff / line_dominant_size * 100)))
                            self._set_run_baseline(run, offset_pct)
                            run.font.size = Pt(line_dominant_size * 0.75)
                        elif bl_diff > line_dominant_size * 0.08:
                            # Subscript (below anchor → negative baseline)
                            offset_pct = max(20, min(50, int(bl_diff / line_dominant_size * 100)))
                            self._set_run_baseline(run, -offset_pct)
                            run.font.size = Pt(line_dominant_size * 0.75)

                # Color
                hex_color = _parse_color(seg.fill_color)
                if hex_color:
                    run.font.color.rgb = RGBColor.from_string(hex_color)

        # No auto-size for paragraphs - use the computed bounding box
        tf.auto_size = MSO_AUTO_SIZE.NONE

        # Small margins for better appearance
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

    def _merge_aligned_equation_clusters(self, clusters: list) -> list:
        """Merge display math clusters that form multi-row aligned equations.

        After initial clustering, aligned equations like:
            f(x) = x² + 2x + 1
                 = (x + 1)²
        may be in separate clusters because they have a vertical gap between
        rows. This method detects such cases by checking:
        1. Overlapping x-ranges (alignment points like '=')
        2. Similar dominant font sizes
        3. Consistent vertical spacing within equation-like range
        4. No text lines between the clusters
        """
        if len(clusters) <= 1:
            return clusters

        def get_baseline(seg):
            return seg.y + seg.font_size * 0.8

        def cluster_bbox(cluster):
            x_min = min(s.x for s in cluster)
            x_max = max(s.x + s.width for s in cluster)
            y_min = min(s.y for s in cluster)
            y_max = max(s.y + s.height for s in cluster)
            return x_min, x_max, y_min, y_max

        def cluster_dominant_size(cluster):
            sizes = [round(s.font_size, 1) for s in cluster]
            from collections import Counter
            return Counter(sizes).most_common(1)[0][0]

        # Sort clusters by y-position
        clusters = sorted(clusters, key=lambda c: min(s.y for s in c))

        merged = True
        while merged:
            merged = False
            new_clusters = []
            used = set()

            for ci in range(len(clusters)):
                if ci in used:
                    continue
                current = list(clusters[ci])
                c_x_min, c_x_max, c_y_min, c_y_max = cluster_bbox(current)
                c_dom = cluster_dominant_size(current)

                for cj in range(ci + 1, len(clusters)):
                    if cj in used:
                        continue
                    other = clusters[cj]
                    o_x_min, o_x_max, o_y_min, o_y_max = cluster_bbox(other)
                    o_dom = cluster_dominant_size(other)

                    # Must have similar dominant font sizes
                    if abs(c_dom - o_dom) > c_dom * 0.2:
                        continue

                    # X-ranges must overlap significantly
                    x_overlap = min(c_x_max, o_x_max) - max(c_x_min, o_x_min)
                    min_width = min(c_x_max - c_x_min, o_x_max - o_x_min)
                    if min_width <= 0 or x_overlap < min_width * 0.3:
                        continue

                    # Y-gap must be reasonable (< 1.5× dominant size)
                    y_gap = max(0, max(o_y_min - c_y_max, c_y_min - o_y_max))
                    if y_gap > c_dom * 1.5:
                        continue

                    # Check for alignment: both should have a '=' segment
                    # at similar x-position (alignment mark)
                    c_eq_xs = [s.x for s in current if s.text.strip() == '=']
                    o_eq_xs = [s.x for s in other if s.text.strip() == '=']
                    has_alignment = False
                    if c_eq_xs and o_eq_xs:
                        for cx in c_eq_xs:
                            for ox in o_eq_xs:
                                if abs(cx - ox) < c_dom * 0.5:
                                    has_alignment = True
                                    break
                            if has_alignment:
                                break

                    # Merge if there's alignment, or if gap is small and
                    # they look like continuation rows
                    if has_alignment or (y_gap < c_dom * 0.8 and x_overlap > min_width * 0.5):
                        current.extend(other)
                        c_x_min, c_x_max, c_y_min, c_y_max = cluster_bbox(current)
                        used.add(cj)
                        merged = True

                new_clusters.append(current)
                used.add(ci)

            clusters = new_clusters

        return clusters

    def _render_math_groups(self, math_segments: list, text_segments: list, slide):
        """Render display math segments grouped spatially into formula regions.

        Math formulas have characters at varying y positions (subscripts,
        superscripts, fractions), so simple baseline grouping doesn't work.

        Strategy:
        1. Cluster math segments into formula groups
        2. Merge clusters that form multi-row aligned equations
        3. Sort segments in reading order using x-position as primary key
        4. Detect sub/superscripts by comparing font size and y-position
           to neighboring anchor segments
        5. Create textbox with per-run font size and baseline offset
        """
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_AUTO_SIZE
        from lxml import etree

        if not math_segments:
            return

        # Group math segments into formula clusters
        math_clusters = self._cluster_math_segments(math_segments)

        # Merge clusters that form aligned multi-row equations
        math_clusters = self._merge_aligned_equation_clusters(math_clusters)

        for cluster in math_clusters:
            font_sizes = [s.font_size for s in cluster]
            # Use most frequent font size as dominant
            size_counts = {}
            for s in font_sizes:
                rs = round(s, 1)
                size_counts[rs] = size_counts.get(rs, 0) + 1
            dominant_size = max(size_counts, key=size_counts.get) if size_counts else 24.8

            def get_baseline(seg):
                return seg.y + seg.font_size * 0.8

            # Deduplicate stacked bracket characters
            # Typst renders tall brackets (for matrices, etc.) as multiple
            # copies of the same bracket character stacked vertically at the
            # same x-position. Keep the one closest to the vertical center
            # of the cluster.
            _BRACKET_CHARS = set('()[]{}⟨⟩⌈⌉⌊⌋|‖')
            bracket_groups = {}  # (x_rounded, char) -> list of segments
            non_bracket_segs = []
            for seg in cluster:
                ch = seg.text.strip()
                if ch in _BRACKET_CHARS:
                    key = (round(seg.x, 0), ch)
                    bracket_groups.setdefault(key, []).append(seg)
                else:
                    non_bracket_segs.append(seg)

            # For each bracket group, keep only the one nearest the vertical center
            cluster_y_center = (min(s.y for s in cluster) +
                                max(s.y + s.height for s in cluster)) / 2
            deduped = list(non_bracket_segs)
            for key, bracket_segs in bracket_groups.items():
                if len(bracket_segs) == 1:
                    deduped.append(bracket_segs[0])
                else:
                    # Pick the one closest to the cluster center
                    best = min(bracket_segs,
                               key=lambda s: abs(s.y + s.height / 2 - cluster_y_center))
                    deduped.append(best)
            cluster = deduped

            # Sort all segments by reading order:
            # Primary: x-position (left to right)
            # Secondary: for segments at similar x, larger font first
            cluster.sort(key=lambda s: (round(s.x, 0), -s.font_size, s.y))

            # Determine anchor baseline for the cluster
            anchor_segs = [s for s in cluster
                           if abs(s.font_size - dominant_size) < dominant_size * 0.25]
            if anchor_segs:
                anchor_bl = sum(get_baseline(s) for s in anchor_segs) / len(anchor_segs)
            else:
                anchor_bl = sum(get_baseline(s) for s in cluster) / len(cluster)

            # Detect multi-row formulas: group segments by baseline bands
            # For display math, detect "anchor rows" where dominant-size segments
            # share a baseline. Each distinct baseline of anchor-size segments
            # is a separate row.
            anchor_only = [s for s in cluster
                           if abs(s.font_size - dominant_size) < dominant_size * 0.25]
            if anchor_only:
                anchor_baselines = sorted(set(round(get_baseline(s), 0) for s in anchor_only))
            else:
                anchor_baselines = sorted(set(round(get_baseline(s), 0) for s in cluster))

            row_bands = []
            if anchor_baselines:
                current_band = [anchor_baselines[0]]
                for bl in anchor_baselines[1:]:
                    # Split into new row if baseline difference > 0.8× dominant size
                    # (standard line spacing for math rows)
                    if bl - current_band[-1] > dominant_size * 0.8:
                        row_bands.append(current_band)
                        current_band = [bl]
                    else:
                        current_band.append(bl)
                row_bands.append(current_band)

            # Multi-row if there are 2+ distinct row bands
            is_multi_row = len(row_bands) > 1

            if is_multi_row:
                # Multi-row formula (e.g., aligned equations with \\)
                row_segments = []
                for band in row_bands:
                    band_min = min(band) - dominant_size * 0.8
                    band_max = max(band) + dominant_size * 0.8
                    row_segs = [s for s in cluster
                                if band_min <= get_baseline(s) <= band_max]
                    row_segs.sort(key=lambda s: (round(s.x, 0), -s.font_size, s.y))
                    if row_segs:
                        row_segments.append(row_segs)
            else:
                row_segments = [sorted(cluster, key=lambda s: (round(s.x, 0), -s.font_size, s.y))]

            # Compute bounding box
            bb_x = min(s.x for s in cluster)
            bb_y = min(s.y for s in cluster)
            bb_right = max(s.x + s.width for s in cluster)
            bb_bottom = max(s.y + s.height for s in cluster)
            bb_w = bb_right - bb_x
            bb_h = bb_bottom - bb_y

            if bb_w < 1:
                bb_w = sum(len(s.text) for s in cluster) * dominant_size * 0.6
            if bb_h < 1:
                bb_h = dominant_size * 1.5 * len(row_segments)

            # Create textbox
            left = Emu(max(0, int(bb_x * EMU_PER_PX)))
            top = Emu(max(0, int(bb_y * EMU_PER_PX)))
            width = Emu(max(int(bb_w * EMU_PER_PX), 9525))
            height = Emu(max(int(bb_h * EMU_PER_PX), 9525))

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = False

            # Render each row as a separate paragraph
            for ri, row_segs in enumerate(row_segments):
                if not row_segs:
                    continue

                if ri == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                # Get the anchor baseline for this row
                row_anchors = [s for s in row_segs
                               if abs(s.font_size - dominant_size) < dominant_size * 0.25]
                if row_anchors:
                    row_anchor_bl = sum(get_baseline(s) for s in row_anchors) / len(row_anchors)
                else:
                    row_anchor_bl = sum(get_baseline(s) for s in row_segs) / len(row_segs)

                prev_seg = None
                for seg in row_segs:
                    run = p.add_run()
                    text = seg.text

                    # Insert space for horizontal gaps
                    if prev_seg is not None:
                        gap = seg.x - (prev_seg.x + prev_seg.width)
                        if gap > dominant_size * 0.15:
                            text = " " + text

                    run.text = text

                    # Use Cambria Math font
                    run.font.name = 'Cambria Math'

                    # Set font size based on the segment's actual size
                    font_size_pt = seg.font_size * 0.75
                    if font_size_pt > 0:
                        run.font.size = Pt(font_size_pt)

                    # Detect superscript/subscript based on font size and baseline offset
                    # PPTX baseline already visually shrinks sup/sub text,
                    # so we use the dominant (parent) font size.
                    size_ratio = seg.font_size / dominant_size if dominant_size > 0 else 1
                    if size_ratio < 0.85:
                        seg_bl = get_baseline(seg)
                        bl_diff = seg_bl - row_anchor_bl

                        if bl_diff < -dominant_size * 0.08:
                            # Superscript (baseline above anchor → positive baseline in PPTX)
                            offset_pct = max(20, min(50, int(-bl_diff / dominant_size * 100)))
                            self._set_run_baseline(run, offset_pct)  # positive = raised
                            # Use dominant font size — PPTX baseline handles shrinking
                            run.font.size = Pt(dominant_size * 0.75)
                        elif bl_diff > dominant_size * 0.08:
                            # Subscript (baseline below anchor → negative baseline in PPTX)
                            offset_pct = max(20, min(50, int(bl_diff / dominant_size * 100)))
                            self._set_run_baseline(run, -offset_pct)  # negative = lowered
                            # Use dominant font size — PPTX baseline handles shrinking
                            run.font.size = Pt(dominant_size * 0.75)

                    # Color
                    hex_color = _parse_color(seg.fill_color)
                    if hex_color:
                        run.font.color.rgb = RGBColor.from_string(hex_color)

                    prev_seg = seg

            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)

    def _render_math_as_curves(self, math_segments: list, slide):
        """Render display math segments as native DrawingML curves (glyph outlines).

        Instead of converting math to text runs with Cambria Math, this method
        renders the original SVG glyph paths as DrawingML custom geometry shapes.
        This ensures visually correct display for complex formulas like matrices,
        integrals, fractions, etc.

        Each math segment has glyph_uses (list of GlyphUse) and glyph_scale.
        The glyph path coordinates are in font units; we apply:
          slide_x = text_x + use.x_offset * glyph_scale
          slide_y = text_y  (baseline position, Y-flipped via scale(S, -S))

        For the path itself, glyph coordinates use Y-up convention.
        We flip Y by negating all Y coords and applying glyph_scale.
        """
        from lxml import etree

        if not math_segments:
            return

        try:
            from ..scripts.svg_to_shapes import (
                parse_svg_path, svg_path_to_absolute,
                normalize_path_commands, path_commands_to_drawingml,
            )
        except ImportError:
            if self.config.verbose:
                print("Warning: svg_to_shapes not available, falling back to text")
            self._render_math_groups(math_segments, [], slide)
            return

        # Cluster and merge math segments (same as text-based rendering)
        math_clusters = self._cluster_math_segments(math_segments)
        math_clusters = self._merge_aligned_equation_clusters(math_clusters)

        for cluster in math_clusters:
            # Collect all glyph shapes for this cluster
            glyph_shapes_xml = []

            for seg in cluster:
                if not seg.glyph_uses or seg.glyph_scale <= 0:
                    continue

                scale = seg.glyph_scale  # e.g., 0.025

                # The segment's position is the top-left of the text box
                # But we need the baseline position for glyph rendering.
                # Recall: seg.y = baseline_y - font_size * 0.8
                # So: baseline_y = seg.y + seg.font_size * 0.8
                baseline_y = seg.y + seg.font_size * 0.8
                seg_x = seg.x

                # Fill color
                fill_color = seg.fill_color or '#000000'
                hex_color = _parse_color(fill_color)
                if not hex_color:
                    hex_color = '000000'

                for glyph_use in seg.glyph_uses:
                    path_d = glyph_use.path_data
                    if not path_d:
                        continue

                    # Parse the glyph path
                    commands = parse_svg_path(path_d)
                    if not commands:
                        continue
                    commands = svg_path_to_absolute(commands)
                    commands = normalize_path_commands(commands)
                    if not commands:
                        continue

                    # The glyph position on the slide:
                    # x = seg_x + glyph_use.x_offset * scale
                    # y = baseline_y + glyph_use.y_offset * (-scale)
                    # Note: y_offset is in font coords (Y-up), so we apply -scale for Y-flip
                    glyph_x = seg_x + glyph_use.x_offset * scale
                    glyph_y = baseline_y + glyph_use.y_offset * (-scale)

                    # Convert path commands to DrawingML
                    # The glyph paths are in font coordinates (Y-up).
                    # We need scale(scale, -scale) to convert to SVG coordinates:
                    #   x_svg = x_font * scale + glyph_x
                    #   y_svg = -y_font * scale + glyph_y (Y-flip)
                    path_xml, min_x, min_y, width, height = path_commands_to_drawingml(
                        commands,
                        offset_x=glyph_x,
                        offset_y=glyph_y,
                        scale_x=scale,
                        scale_y=-scale,  # Y-flip for font coordinates
                    )

                    if not path_xml or width <= 0 or height <= 0:
                        continue

                    # Build the individual glyph shape XML
                    w_emu = int(width * EMU_PER_PX)
                    h_emu = int(height * EMU_PER_PX)
                    off_x_emu = int(min_x * EMU_PER_PX)
                    off_y_emu = int(min_y * EMU_PER_PX)

                    shape_id = len(slide.shapes) + 200 + len(glyph_shapes_xml)

                    shape_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
<p:nvSpPr>
<p:cNvPr id="{shape_id}" name="MathGlyph {shape_id}"/>
<p:cNvSpPr/><p:nvPr/>
</p:nvSpPr>
<p:spPr>
<a:xfrm><a:off x="{off_x_emu}" y="{off_y_emu}"/><a:ext cx="{w_emu}" cy="{h_emu}"/></a:xfrm>
<a:custGeom>
<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
<a:rect l="l" t="t" r="r" b="b"/>
<a:pathLst><a:path w="{w_emu}" h="{h_emu}">
{path_xml}
</a:path></a:pathLst>
</a:custGeom>
<a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
<a:ln w="0"><a:noFill/></a:ln>
</p:spPr>
</p:sp>'''

                    glyph_shapes_xml.append(shape_xml)

            # Inject glyph shapes into the slide, grouped per cluster
            if not glyph_shapes_xml:
                continue

            sp_tree = slide.shapes._spTree

            if len(glyph_shapes_xml) > 1:
                # Wrap in a group shape (<p:grpSp>) so user can drag the formula as one unit
                try:
                    # Parse all individual shapes first to compute bounding box
                    shape_elems = []
                    all_x = []
                    all_y = []
                    all_r = []
                    all_b = []
                    for sxml in glyph_shapes_xml:
                        sel = etree.fromstring(sxml)
                        shape_elems.append(sel)
                        # Extract position from <a:off x="..." y="..."/>
                        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                        xfrm = sel.find(f'.//{ns_a}xfrm')
                        if xfrm is not None:
                            off = xfrm.find(f'{ns_a}off')
                            ext = xfrm.find(f'{ns_a}ext')
                            if off is not None and ext is not None:
                                ox = int(off.get('x', '0'))
                                oy = int(off.get('y', '0'))
                                cx = int(ext.get('cx', '0'))
                                cy = int(ext.get('cy', '0'))
                                all_x.append(ox)
                                all_y.append(oy)
                                all_r.append(ox + cx)
                                all_b.append(oy + cy)

                    if all_x:
                        grp_x = min(all_x)
                        grp_y = min(all_y)
                        grp_cx = max(all_r) - grp_x
                        grp_cy = max(all_b) - grp_y
                    else:
                        grp_x = grp_y = 0
                        grp_cx = grp_cy = EMU_PER_PX

                    grp_id = len(slide.shapes) + 300 + len(sp_tree)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

                    grp_sp = etree.SubElement(sp_tree, f'{{{ns_p}}}grpSp')
                    # nvGrpSpPr
                    nv = etree.SubElement(grp_sp, f'{{{ns_p}}}nvGrpSpPr')
                    cnvpr = etree.SubElement(nv, f'{{{ns_p}}}cNvPr')
                    cnvpr.set('id', str(grp_id))
                    cnvpr.set('name', f'MathFormula {grp_id}')
                    etree.SubElement(nv, f'{{{ns_p}}}cNvGrpSpPr')
                    etree.SubElement(nv, f'{{{ns_p}}}nvPr')
                    # grpSpPr with transform
                    grp_sp_pr = etree.SubElement(grp_sp, f'{{{ns_p}}}grpSpPr')
                    xfrm = etree.SubElement(grp_sp_pr, f'{{{ns_a}}}xfrm')
                    off = etree.SubElement(xfrm, f'{{{ns_a}}}off')
                    off.set('x', str(grp_x))
                    off.set('y', str(grp_y))
                    ext = etree.SubElement(xfrm, f'{{{ns_a}}}ext')
                    ext.set('cx', str(grp_cx))
                    ext.set('cy', str(grp_cy))
                    ch_off = etree.SubElement(xfrm, f'{{{ns_a}}}chOff')
                    ch_off.set('x', str(grp_x))
                    ch_off.set('y', str(grp_y))
                    ch_ext = etree.SubElement(xfrm, f'{{{ns_a}}}chExt')
                    ch_ext.set('cx', str(grp_cx))
                    ch_ext.set('cy', str(grp_cy))

                    # Append each glyph shape as child of group
                    for sel in shape_elems:
                        grp_sp.append(sel)
                except Exception as e:
                    if self.config.verbose:
                        print(f"Warning: Failed to create math group shape: {e}")
                    # Fallback: add individually
                    for sxml in glyph_shapes_xml:
                        try:
                            sp_tree.append(etree.fromstring(sxml))
                        except Exception:
                            pass
            else:
                # Single glyph, no need for grouping
                try:
                    sp_tree.append(etree.fromstring(glyph_shapes_xml[0]))
                except Exception as e:
                    if self.config.verbose:
                        print(f"Warning: Failed to inject math glyph shape: {e}")

    def _set_run_baseline(self, run, baseline_pct: int):
        """Set the baseline offset for a run (superscript/subscript).

        Args:
            run: A python-pptx Run object
            baseline_pct: Baseline offset as a percentage (OOXML convention).
                         Positive = superscript (text raised up).
                         Negative = subscript (text lowered down).
                         E.g., 30 for superscript, -30 for subscript.
        """
        from lxml import etree

        # Access the run's rPr element
        rPr = run._r.get_or_add_rPr()
        # Set baseline attribute (in 1000ths of a percent)
        baseline_val = baseline_pct * 1000
        rPr.set('baseline', str(baseline_val))

    def _cluster_math_segments(self, math_segments: list) -> List[List]:
        """Cluster math segments into formula groups.

        Uses a two-phase approach:
        1. Proximity-based clustering: Group all segments (not just anchors)
           by spatial proximity using both x and y distance
        2. Split by large horizontal gaps within each cluster

        This handles fractions, radicals, and matrices better than the previous
        anchor-only approach, because all parts of a formula (numerator,
        denominator, operators) get grouped together by spatial proximity.
        """
        if not math_segments:
            return []

        font_sizes = [s.font_size for s in math_segments]
        # Use the most frequent font size as the "dominant" size
        size_counts = {}
        for s in font_sizes:
            rs = round(s, 1)
            size_counts[rs] = size_counts.get(rs, 0) + 1
        dominant_size = max(size_counts, key=size_counts.get) if size_counts else 24.8

        def get_baseline(seg):
            return seg.y + seg.font_size * 0.8

        def seg_center_y(seg):
            return seg.y + seg.height / 2

        # Phase 1: Sort all segments by x-position, then group by spatial proximity
        # This is a union-find/connected-component approach
        all_segs = list(math_segments)
        n = len(all_segs)

        if n == 1:
            return [all_segs]

        # Compute bounding boxes for proximity checks
        # Two segments belong to the same formula if:
        # - They overlap in x-range (or are close) AND y-range overlap/proximity
        # - OR they are horizontally adjacent at similar y position

        # Use union-find to group segments
        parent = list(range(n))

        def find(i):
            while parent[i] != i:
                parent[i] = parent[parent[i]]
                i = parent[i]
            return i

        def union(i, j):
            ri, rj = find(i), find(j)
            if ri != rj:
                parent[ri] = rj

        # For each pair of segments, check if they should be in the same cluster
        # Optimization: sort by x and only compare nearby segments
        sorted_indices = sorted(range(n), key=lambda i: all_segs[i].x)

        for ii in range(len(sorted_indices)):
            i = sorted_indices[ii]
            si = all_segs[i]
            si_right = si.x + si.width
            si_bl = get_baseline(si)

            for jj in range(ii + 1, len(sorted_indices)):
                j = sorted_indices[jj]
                sj = all_segs[j]

                # Stop looking if segments are too far apart in x
                x_gap = sj.x - si_right
                if x_gap > dominant_size * 2.5:
                    break

                sj_bl = get_baseline(sj)
                bl_diff = abs(si_bl - sj_bl)

                # Check various proximity conditions for grouping

                # Condition 1: Same baseline and close in x
                # (standard reading order - adjacent characters on same line)
                if bl_diff < dominant_size * 0.5 and x_gap < dominant_size * 1.5:
                    union(i, j)
                    continue

                # Condition 2: Superscript/subscript detection
                # One segment is smaller and positioned right after/above/below the other
                size_ratio = min(si.font_size, sj.font_size) / max(si.font_size, sj.font_size)
                if size_ratio < 0.85:
                    # Size difference suggests sub/superscript relationship
                    # Must be horizontally close (touching or slightly overlapping)
                    if x_gap < dominant_size * 0.5:
                        if bl_diff < dominant_size * 1.2:
                            union(i, j)
                            continue

                # Condition 3: Fraction numerator/denominator detection
                # Two segments at the SAME x-position but different y positions
                # Must have significant x-overlap AND be at the same font size
                # (fractions have similar-sized numerator and denominator)
                si_x_min, si_x_max = si.x, si.x + si.width
                sj_x_min, sj_x_max = sj.x, sj.x + sj.width
                x_overlap = min(si_x_max, sj_x_max) - max(si_x_min, sj_x_min)
                min_width = min(si.width, sj.width)

                if min_width > 0 and x_overlap > min_width * 0.5:
                    # Significant x-overlap: possible fraction parts
                    # Check that they are the same font size (fraction parts)
                    # and vertically close (within 1.2× dominant size)
                    same_size = abs(si.font_size - sj.font_size) < dominant_size * 0.15
                    if same_size and bl_diff < dominant_size * 1.2:
                        union(i, j)
                        continue

        # Collect clusters from union-find
        groups = {}
        for i in range(n):
            root = find(i)
            groups.setdefault(root, []).append(all_segs[i])

        clusters = list(groups.values())

        # Phase 2: Merge clusters that are horizontally touching/overlapping
        # and vertically overlapping. This catches formula parts separated
        # by the initial clustering (like x= + fraction, √ + radicand).
        # Criteria: x-ranges must overlap or be very close (touching),
        # AND y-ranges must overlap significantly.
        merged = True
        while merged:
            merged = False
            new_clusters = []
            used = set()
            for ci in range(len(clusters)):
                if ci in used:
                    continue
                current = list(clusters[ci])
                c_x_min = min(s.x for s in current)
                c_x_max = max(s.x + s.width for s in current)
                c_y_min = min(s.y for s in current)
                c_y_max = max(s.y + s.height for s in current)

                for cj in range(ci + 1, len(clusters)):
                    if cj in used:
                        continue
                    other = clusters[cj]
                    o_x_min = min(s.x for s in other)
                    o_x_max = max(s.x + s.width for s in other)
                    o_y_min = min(s.y for s in other)
                    o_y_max = max(s.y + s.height for s in other)

                    # X-ranges must be touching or overlapping
                    x_gap = max(0, max(o_x_min - c_x_max, c_x_min - o_x_max))
                    # Y-ranges must overlap (positive overlap)
                    y_overlap = min(c_y_max, o_y_max) - max(c_y_min, o_y_min)

                    should_merge = False

                    # Case 1: X-ranges overlap AND y-ranges overlap or nearly touch
                    # (fraction parts, radical + radicand, matrix rows, etc.)
                    # Use small tolerance (1px) for y-overlap to handle boundary cases
                    if x_gap == 0 and y_overlap >= -1.0:
                        should_merge = True

                    # Case 2: Very close horizontally (< 0.5× dom size)
                    # AND significant y-overlap (> 30% of smaller cluster height)
                    elif x_gap < dominant_size * 0.5 and y_overlap > 0:
                        smaller_h = min(c_y_max - c_y_min, o_y_max - o_y_min)
                        if smaller_h > 0 and y_overlap > smaller_h * 0.3:
                            should_merge = True

                    if should_merge:
                        current.extend(other)
                        c_x_min = min(c_x_min, o_x_min)
                        c_x_max = max(c_x_max, o_x_max)
                        c_y_min = min(c_y_min, o_y_min)
                        c_y_max = max(c_y_max, o_y_max)
                        used.add(cj)
                        merged = True

                new_clusters.append(current)
                used.add(ci)
            clusters = new_clusters

        # Phase 3: Within each cluster, split by very large x-gaps
        # This handles cases where two separate formulas got merged
        final_clusters = []
        for cluster in clusters:
            if len(cluster) <= 1:
                final_clusters.append(cluster)
                continue

            cluster.sort(key=lambda s: s.x)

            # Find dominant font size for this cluster
            cl_sizes = [round(s.font_size, 1) for s in cluster]
            cl_size_counts = {}
            for s in cl_sizes:
                cl_size_counts[s] = cl_size_counts.get(s, 0) + 1
            cl_dominant = max(cl_size_counts, key=cl_size_counts.get)

            current = [cluster[0]]
            for seg in cluster[1:]:
                prev = current[-1]
                gap = seg.x - (prev.x + prev.width)
                # Large gap: separate formulas
                if gap > cl_dominant * 4.0:
                    final_clusters.append(current)
                    current = [seg]
                else:
                    current.append(seg)
            final_clusters.append(current)

        return final_clusters


def _find_typst_ts_cli() -> str:
    """Locate the typst-ts-cli binary.

    Search order:
    1. Bundled binary shipped inside this package (typ2pptx/data/bin/)
    2. System PATH
    """
    # 1. Bundled binary
    bundled = Path(__file__).parent.parent / "data" / "bin" / "typst-ts-cli"
    if bundled.exists() and os.access(str(bundled), os.X_OK):
        return str(bundled)

    # 2. System PATH
    system_cli = shutil.which("typst-ts-cli")
    if system_cli:
        return system_cli

    raise FileNotFoundError(
        "typst-ts-cli not found. Install it from "
        "https://github.com/Myriad-Dreamin/typst.ts/releases "
        "or specify its path with --typst-ts-cli"
    )


def compile_typst_to_svg(
    typ_path: str,
    output_svg: Optional[str] = None,
    typst_ts_cli: Optional[str] = None,
    root: Optional[str] = None,
) -> str:
    """Compile a .typ file to SVG using typst-ts-cli.

    Args:
        typ_path: Path to the .typ file
        output_svg: Output SVG path (default: same name with .artifact.svg)
        typst_ts_cli: Path to the typst-ts-cli binary (default: bundled or system)
        root: Root directory for the Typst project (for resolving imports/paths)

    Returns:
        Path to the generated SVG file
    """
    typ_path = Path(typ_path).resolve()
    if not typ_path.exists():
        raise FileNotFoundError(f"Typst file not found: {typ_path}")

    if typst_ts_cli is None:
        typst_ts_cli = _find_typst_ts_cli()

    workspace_dir = root if root else str(typ_path.parent)

    cmd = [
        typst_ts_cli,
        "compile",
        "--workspace", workspace_dir,
        "--entry", str(typ_path),
        "--format", "svg",
    ]

    if output_svg:
        cmd.extend(["-o", output_svg])

    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        cwd=workspace_dir,
    )

    if result.returncode != 0:
        raise RuntimeError(
            f"typst-ts-cli failed:\nstdout: {result.stdout}\nstderr: {result.stderr}"
        )

    # Find the output file
    if output_svg:
        return output_svg
    else:
        # Default output is {name}.artifact.svg
        default_output = typ_path.with_suffix('.artifact.svg')
        if default_output.exists():
            return str(default_output)
        # Try without .artifact
        alt_output = typ_path.with_suffix('.svg')
        if alt_output.exists():
            return str(alt_output)
        raise FileNotFoundError(
            f"Could not find SVG output. Expected: {default_output}"
        )


def query_speaker_notes(
    typ_path: str,
    root: Optional[str] = None,
) -> Dict[int, str]:
    """Query speaker notes from a .typ file using the typst Python package.

    Uses typst.query() to extract pdfpc-format speaker notes from Touying
    presentations.

    Args:
        typ_path: Path to the .typ file
        root: Root directory for the Typst project (for resolving imports/paths)

    Returns:
        Dict mapping page index (0-based) to notes text
    """
    try:
        import typst as typst_py

        typ_path_resolved = str(Path(typ_path).resolve())
        query_kwargs: Dict[str, Any] = {"field": "value"}
        if root is not None:
            query_kwargs["root"] = root
        result = typst_py.query(typ_path_resolved, "<pdfpc-file>", **query_kwargs)

        pdfpc = json.loads(result)
        if not pdfpc:
            return {}

        pdfpc = pdfpc[0] if isinstance(pdfpc, list) and pdfpc else pdfpc
        if not pdfpc or 'pages' not in pdfpc:
            return {}

        return {
            page['idx']: page['note']
            for page in pdfpc['pages']
            if 'note' in page
        }

    except Exception:
        return {}


def convert_typst_to_pptx(
    input_path: str,
    output_path: Optional[str] = None,
    typst_ts_cli: Optional[str] = None,
    root: Optional[str] = None,
    config: Optional[ConversionConfig] = None,
    verbose: bool = False,
) -> str:
    """Convert a Typst file or SVG to PPTX.

    Args:
        input_path: Path to .typ or .svg file
        output_path: Output .pptx path (default: same name)
        typst_ts_cli: Path to typst-ts-cli binary (default: bundled or system)
        root: Root directory for the Typst project (for resolving imports/paths)
        config: Conversion configuration
        verbose: Enable verbose output

    Returns:
        Path to the generated PPTX file
    """
    config = config or ConversionConfig()
    config.verbose = verbose

    input_path = Path(input_path)
    if output_path is None:
        output_path = str(input_path.with_suffix('.pptx'))

    svg_path = None

    if input_path.suffix == '.typ':
        # Compile to SVG first
        if verbose:
            print(f"Compiling {input_path} to SVG...")
        svg_path = compile_typst_to_svg(
            str(input_path), typst_ts_cli=typst_ts_cli, root=root,
        )

        # Query speaker notes
        speaker_notes = query_speaker_notes(str(input_path), root=root)
    elif input_path.suffix == '.svg':
        svg_path = str(input_path)
        speaker_notes = {}
    else:
        raise ValueError(f"Unsupported input format: {input_path.suffix}")

    if verbose:
        print(f"Converting {svg_path} to PPTX...")

    converter = TypstSVGConverter(config)
    result = converter.convert(svg_path, output_path, speaker_notes=speaker_notes)

    if verbose:
        print(f"Done! Output: {result}")

    return result
