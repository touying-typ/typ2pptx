"""Tests for math formula detection and rendering."""
import pytest
from pptx import Presentation
from pptx.util import Emu

from typst2pptx.core.typst_svg_parser import parse_typst_svg


class TestMathFontDetection:
    """Test that math fonts are correctly detected."""

    def test_math_font_detected(self, math_test_parsed):
        """Math font variant should be detected as 'math'."""
        styles = {v.style for v in math_test_parsed.font_variants.values()}
        assert 'math' in styles, f"No math font variant found, styles: {styles}"

    def test_regular_font_not_misclassified(self, math_test_parsed):
        """Regular text font should not be classified as math or italic."""
        styles = {v.style for v in math_test_parsed.font_variants.values()}
        assert 'regular' in styles, f"No regular font variant found, styles: {styles}"

    def test_bold_font_detected(self, math_test_parsed):
        """Bold (heading) font should still be detected correctly."""
        styles = {v.style for v in math_test_parsed.font_variants.values()}
        assert 'bold' in styles, f"No bold font variant found, styles: {styles}"

    def test_math_segments_have_math_variant(self, math_test_parsed):
        """Math text segments should have font_variant='math'."""
        math_chars = set('𝑒𝑖𝜋𝑥𝑏𝑎𝑐𝑛')
        page1 = math_test_parsed.pages[0]
        for seg in page1.text_segments:
            if any(c in seg.text for c in math_chars):
                assert seg.font_variant == 'math', (
                    f"Math segment '{seg.text}' has variant '{seg.font_variant}'"
                )

    def test_regular_text_not_math(self, math_test_parsed):
        """Regular body text should not be classified as math."""
        page1 = math_test_parsed.pages[0]
        for seg in page1.text_segments:
            if seg.text.startswith("The "):
                assert seg.font_variant == 'regular', (
                    f"Regular text '{seg.text}' has variant '{seg.font_variant}'"
                )


class TestMathPPTXOutput:
    """Test the math formula PPTX output."""

    def test_math_slide_count(self, math_test_pptx):
        """math_test.typ should produce 3 slides."""
        prs = Presentation(math_test_pptx)
        assert len(prs.slides) == 3

    def test_math_has_cambria_font(self, math_test_pptx):
        """Math formulas should use Cambria Math font."""
        prs = Presentation(math_test_pptx)
        has_cambria = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            if run.font.name == 'Cambria Math':
                                has_cambria = True
        assert has_cambria, "No Cambria Math font found in PPTX"

    def test_regular_text_uses_arial(self, math_test_pptx):
        """Regular text should use Arial font."""
        prs = Presentation(math_test_pptx)
        has_arial = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            if run.font.name == 'Arial' and len(run.text) > 5:
                                has_arial = True
        assert has_arial, "No Arial font for regular text found in PPTX"

    def test_math_formula_text_readable(self, math_test_pptx):
        """Math formula text should contain recognizable math symbols."""
        prs = Presentation(math_test_pptx)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            if run.font.name == 'Cambria Math':
                                all_text.append(run.text)

        combined = " ".join(all_text)
        # Should contain math symbols
        assert any(c in combined for c in '𝑒𝜋𝑥+=∫∑√'), (
            f"Math text missing expected symbols: '{combined[:100]}'"
        )

    def test_math_titles_are_bold(self, math_test_pptx):
        """Slide titles should be bold."""
        prs = Presentation(math_test_pptx)
        # Check slide 1 title
        slide1 = prs.slides[0]
        found_bold_title = False
        for shape in slide1.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if "Math Formulas" in run.text and run.font.bold:
                            found_bold_title = True
        assert found_bold_title, "Title 'Math Formulas' should be bold"

    def test_inline_math_with_text(self, math_test_pptx):
        """Inline math should appear in the same textbox as surrounding text,
        or be positioned correctly relative to it."""
        prs = Presentation(math_test_pptx)
        slide1 = prs.slides[0]

        # Check for inline math merged with text in same shape
        found_inline = False
        for shape in slide1.shapes:
            if shape.has_text_frame:
                runs = []
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        runs.append(run)
                all_text = " ".join(r.text for r in runs)
                has_text = "identity" in all_text
                has_math = any(r.font.name == 'Cambria Math' and '𝑒' in r.text for r in runs)
                if has_text and has_math:
                    found_inline = True
                    break

        if not found_inline:
            # Alternatively, check that math is in a separate shape but positioned
            # after the text shape
            text_shapes = []
            math_shapes = []
            for shape in slide1.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            if "identity" in run.text:
                                text_shapes.append(shape)
                            elif run.font.name == 'Cambria Math' and '𝑒' in run.text:
                                math_shapes.append(shape)

            if text_shapes and math_shapes:
                text_right = text_shapes[0].left + text_shapes[0].width
                math_left = math_shapes[0].left
                assert math_left >= text_right - Emu(50000), (
                    "Inline math should be positioned after the preceding text"
                )
                found_inline = True

        assert found_inline, "No inline math found near 'identity' text"

    def test_matrix_slide_has_math(self, math_test_pptx):
        """Matrix slide should have math content (rendered as curves or text)."""
        from lxml import etree
        prs = Presentation(math_test_pptx)
        slide3 = prs.slides[2]

        # Display math is now rendered as native DrawingML glyph curves
        # (MathGlyph shapes) instead of Cambria Math text runs.
        # Check for either MathGlyph shapes or Cambria Math text.
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        sp_tree = slide3.shapes._spTree
        math_glyph_count = 0
        for sp in sp_tree.findall('.//p:sp', ns):
            cnvpr = sp.find('.//p:nvSpPr/p:cNvPr', ns)
            if cnvpr is not None and 'MathGlyph' in cnvpr.get('name', ''):
                math_glyph_count += 1

        math_text = []
        for shape in slide3.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.name == 'Cambria Math':
                            math_text.append(run.text)

        has_math = math_glyph_count > 0 or len(math_text) > 0
        assert has_math, (
            f"Matrix slide has no math content: "
            f"{math_glyph_count} MathGlyph shapes, "
            f"Cambria Math text: '{' '.join(math_text)}'"
        )


class TestMathSVGParsing:
    """Test SVG parsing for math content."""

    def test_math_page_count(self, math_test_parsed):
        """math_test.typ should produce 3 pages."""
        assert len(math_test_parsed.pages) == 3

    def test_math_has_text_segments(self, math_test_parsed):
        """Each page should have text segments."""
        for page in math_test_parsed.pages:
            assert len(page.text_segments) > 0, (
                f"Page {page.page_num} has no text segments"
            )

    def test_math_segments_have_positions(self, math_test_parsed):
        """Math text segments should have valid positions."""
        for page in math_test_parsed.pages:
            for seg in page.text_segments:
                assert isinstance(seg.x, float), f"x is not float: {type(seg.x)}"
                assert isinstance(seg.y, float), f"y is not float: {type(seg.y)}"
                assert seg.font_size > 0, f"Invalid font size: {seg.font_size}"

    def test_math_font_sizes_vary(self, math_test_parsed):
        """Math should have varying font sizes (main, sub/superscript)."""
        page1 = math_test_parsed.pages[0]
        math_segs = [s for s in page1.text_segments if s.font_variant == 'math']
        if math_segs:
            sizes = {round(s.font_size, 1) for s in math_segs}
            assert len(sizes) > 1, (
                f"Expected varying math font sizes, got: {sizes}"
            )
