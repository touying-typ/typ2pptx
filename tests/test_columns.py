"""Tests for two-column layout paragraph merging."""
import pytest
from pptx import Presentation


class TestColumnsSVGParsing:
    """Test SVG parsing for two-column layout content."""

    def test_columns_page_count(self, columns_test_parsed):
        """columns_test.typ should produce at least 2 slides (title + content)."""
        assert len(columns_test_parsed.pages) >= 2

    def test_columns_has_text_segments(self, columns_test_parsed):
        """The content slide should have text segments from both columns."""
        content_page = columns_test_parsed.pages[1]
        assert len(content_page.text_segments) > 5, (
            "Two-column lorem should produce many text segments"
        )

    def test_columns_two_distinct_x_positions(self, columns_test_parsed):
        """Text segments should cluster around two distinct x positions (two columns)."""
        content_page = columns_test_parsed.pages[1]
        body_segs = [
            s for s in content_page.text_segments
            if s.font_size > 20
        ]
        x_positions = sorted(set(round(s.x) for s in body_segs))
        assert len(x_positions) >= 2, (
            f"Expected at least 2 distinct x positions for two columns, got {x_positions}"
        )


class TestColumnsPPTXOutput:
    """Test two-column layout PPTX output."""

    def test_columns_slide_count(self, columns_test_pptx):
        """columns_test.typ should produce slides."""
        prs = Presentation(columns_test_pptx)
        assert len(prs.slides) >= 2

    def _get_column_shapes(self, pptx_path):
        """Helper: get left and right column body text shapes from slide 2."""
        prs = Presentation(pptx_path)
        content_slide = prs.slides[1]
        slide_width = prs.slide_width
        midpoint = slide_width // 2

        left_shapes = []
        right_shapes = []
        for shape in content_slide.shapes:
            if not shape.has_text_frame:
                continue
            all_text = ""
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    all_text += run.text
            # Filter to body text only (skip title, page numbers, etc.)
            if len(all_text.strip()) < 20:
                continue
            if "Column" in all_text or "Lorem" in all_text[:15]:
                # Skip title-like shapes, but keep "Lorem ipsum..." body
                if "Lorem" not in all_text:
                    continue
            if shape.left < midpoint:
                left_shapes.append((shape, all_text))
            else:
                right_shapes.append((shape, all_text))
        return left_shapes, right_shapes

    def test_left_column_not_merged_by_default(self, columns_test_pptx):
        """With the default config (paragraph detection OFF), each wrapped line
        in the left column should remain its own textbox.

        Paragraph auto-detection is disabled by default because the heuristic
        often mis-merges tightly-packed content (tables, bullets, etc.). With
        detection off, the ~10 wrapped lines in the left column should produce
        noticeably more than 2 shapes.
        """
        left_shapes, _ = self._get_column_shapes(columns_test_pptx)

        assert len(left_shapes) > 0, "No left column body text found"
        assert len(left_shapes) >= 3, (
            f"With detect_paragraphs=False (default), left column's wrapped "
            f"lines should stay as separate textboxes, but only found "
            f"{len(left_shapes)} shapes (looks like they were merged)."
        )

    def test_right_column_not_merged_by_default(self, columns_test_pptx):
        """Same as the left-column case: default config keeps each wrapped
        line as its own textbox."""
        _, right_shapes = self._get_column_shapes(columns_test_pptx)

        assert len(right_shapes) > 0, "No right column body text found"
        assert len(right_shapes) >= 3, (
            f"With detect_paragraphs=False (default), right column's wrapped "
            f"lines should stay as separate textboxes, but only found "
            f"{len(right_shapes)} shapes (looks like they were merged)."
        )

    def test_left_column_merged_when_detection_enabled(
        self, columns_test_pptx_with_paragraphs,
    ):
        """With detect_paragraphs=True, the left column's multi-line lorem text
        should be merged into 1-2 paragraph textboxes."""
        left_shapes, _ = self._get_column_shapes(columns_test_pptx_with_paragraphs)

        assert len(left_shapes) > 0, "No left column body text found"
        assert len(left_shapes) <= 2, (
            f"With detect_paragraphs=True, left column text should be merged "
            f"into 1-2 paragraph textboxes, but found {len(left_shapes)} "
            f"separate shapes"
        )

    def test_right_column_merged_when_detection_enabled(
        self, columns_test_pptx_with_paragraphs,
    ):
        """With detect_paragraphs=True, the right column's multi-line lorem
        text should be merged into 1-2 paragraph textboxes."""
        _, right_shapes = self._get_column_shapes(columns_test_pptx_with_paragraphs)

        assert len(right_shapes) > 0, "No right column body text found"
        assert len(right_shapes) <= 2, (
            f"With detect_paragraphs=True, right column text should be merged "
            f"into 1-2 paragraph textboxes, but found {len(right_shapes)} "
            f"separate shapes"
        )

    def test_columns_text_not_overlapping(self, columns_test_pptx):
        """Left and right column textboxes should not overlap horizontally."""
        prs = Presentation(columns_test_pptx)
        content_slide = prs.slides[1]

        body_shapes = []
        for shape in content_slide.shapes:
            if not shape.has_text_frame:
                continue
            all_text = ""
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    all_text += run.text
            if len(all_text) > 20:
                body_shapes.append(shape)

        if len(body_shapes) >= 2:
            lefts = sorted(body_shapes, key=lambda s: s.left)
            leftmost = lefts[0]
            rightmost = lefts[-1]
            # The rightmost shape's left should be beyond the leftmost shape's right edge
            left_right_edge = leftmost.left + leftmost.width
            assert rightmost.left >= left_right_edge * 0.9, (
                f"Columns should not overlap: left column right edge={left_right_edge}, "
                f"right column left={rightmost.left}"
            )
