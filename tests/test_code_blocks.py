"""Tests for code block rendering."""
import pytest
from pptx import Presentation


class TestCodeBlockSVGParsing:
    """Test SVG parsing for code block content."""

    def test_code_block_page_count(self, code_block_test_parsed):
        """code_block_test.typ should produce 5 slides."""
        assert len(code_block_test_parsed.pages) == 5

    def test_inline_code_detected_as_mono(self, code_block_test_parsed):
        """Inline code should be detected as monospace font variant."""
        page2 = code_block_test_parsed.pages[1]
        mono_segs = [s for s in page2.text_segments if s.font_variant == 'mono']
        assert len(mono_segs) >= 3, (
            f"Slide 2 should have at least 3 mono segments (inline code), "
            f"got {len(mono_segs)}"
        )

    def test_code_block_has_mono_segments(self, code_block_test_parsed):
        """Code block slide should have many monospace segments."""
        page3 = code_block_test_parsed.pages[2]
        mono_segs = [s for s in page3.text_segments if s.font_variant == 'mono']
        assert len(mono_segs) >= 10, (
            f"Slide 3 (Python code block) should have many mono segments, "
            f"got {len(mono_segs)}"
        )

    def test_code_block_has_syntax_colors(self, code_block_test_parsed):
        """Code block should have multiple distinct fill colors (syntax highlighting)."""
        page3 = code_block_test_parsed.pages[2]
        mono_segs = [s for s in page3.text_segments if s.font_variant == 'mono']
        colors = set(s.fill_color for s in mono_segs if s.fill_color)
        assert len(colors) >= 2, (
            f"Code block should have at least 2 syntax colors, got {colors}"
        )


class TestCodeBlockPPTXOutput:
    """Test code block PPTX output."""

    def test_code_block_slide_count(self, code_block_test_pptx):
        """code_block_test.typ should produce 5 slides."""
        prs = Presentation(code_block_test_pptx)
        assert len(prs.slides) == 5

    def test_inline_code_uses_consolas(self, code_block_test_pptx):
        """Inline code should use Consolas font."""
        prs = Presentation(code_block_test_pptx)
        slide2 = prs.slides[1]

        consolas_texts = []
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name == 'Consolas':
                            consolas_texts.append(r.text)

        assert any('inline code' in t for t in consolas_texts), (
            f"Expected 'inline code' in Consolas runs, got: {consolas_texts}"
        )

    def test_code_block_uses_consolas(self, code_block_test_pptx):
        """Code block content should use Consolas font."""
        prs = Presentation(code_block_test_pptx)
        slide3 = prs.slides[2]

        consolas_count = 0
        for shape in slide3.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name == 'Consolas':
                            consolas_count += 1

        assert consolas_count >= 10, (
            f"Python code block should have at least 10 Consolas runs, "
            f"got {consolas_count}"
        )

    def test_python_keywords_present(self, code_block_test_pptx):
        """Python code block should contain keywords like 'def', 'for', 'print'."""
        prs = Presentation(code_block_test_pptx)
        slide3 = prs.slides[2]

        all_code_text = []
        for shape in slide3.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name == 'Consolas':
                            all_code_text.append(r.text)

        combined = " ".join(all_code_text)
        assert 'def' in combined, f"'def' not found in code text: {combined[:200]}"
        assert 'hello' in combined, f"'hello' not found in code text: {combined[:200]}"

    def test_syntax_highlighting_colors(self, code_block_test_pptx):
        """Code block should have syntax-highlighted colors (not all black)."""
        prs = Presentation(code_block_test_pptx)
        slide3 = prs.slides[2]

        colors = set()
        for shape in slide3.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name == 'Consolas' and r.font.color.rgb:
                            colors.add(str(r.font.color.rgb))

        assert len(colors) >= 2, (
            f"Code block should have at least 2 distinct syntax colors, "
            f"got {colors}"
        )

    def test_rust_code_block(self, code_block_test_pptx):
        """Rust code block should have recognizable content."""
        prs = Presentation(code_block_test_pptx)
        slide4 = prs.slides[3]

        all_code_text = []
        for shape in slide4.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name == 'Consolas':
                            all_code_text.append(r.text)

        combined = " ".join(all_code_text)
        assert 'fn' in combined or 'main' in combined, (
            f"Rust keywords not found in code text: {combined[:200]}"
        )

    def test_mixed_text_and_code(self, code_block_test_pptx):
        """Slide with mixed content should have both regular and Consolas text."""
        prs = Presentation(code_block_test_pptx)
        slide5 = prs.slides[4]

        has_regular = False
        has_consolas = False
        for shape in slide5.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name == 'Consolas':
                            has_consolas = True
                        elif r.font.name and r.font.name != 'Consolas':
                            has_regular = True

        assert has_regular, "Should have regular (non-code) text"
        assert has_consolas, "Should have Consolas (code) text"

    def test_inline_code_mixed_with_text(self, code_block_test_pptx):
        """Inline code should appear in the same textbox as surrounding text."""
        prs = Presentation(code_block_test_pptx)
        slide2 = prs.slides[1]

        found = False
        for shape in slide2.shapes:
            if shape.has_text_frame:
                has_regular = False
                has_mono = False
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name == 'Consolas':
                            has_mono = True
                        elif r.text.strip():
                            has_regular = True
                if has_regular and has_mono:
                    found = True

        assert found, "Inline code should be in same textbox as regular text"
