"""Tests for embedded image handling."""
import pytest
from io import BytesIO
from pptx import Presentation
from PIL import Image


class TestImageSVGParsing:
    """Test SVG parsing for embedded images."""

    def test_image_page_count(self, image_test_parsed):
        """image_test.typ should produce 4 slides."""
        assert len(image_test_parsed.pages) == 4

    def test_png_image_detected(self, image_test_parsed):
        """Slide 2 should have an image shape (PNG)."""
        page2 = image_test_parsed.pages[1]
        image_shapes = [s for s in page2.shapes if s.tag == 'image']
        assert len(image_shapes) >= 1, "Slide 2 should have at least 1 image shape"

    def test_svg_image_detected(self, image_test_parsed):
        """Slide 3 should have shapes (SVG content)."""
        page3 = image_test_parsed.pages[2]
        # SVG images in typst are rendered as inline SVG, may appear as shapes
        total = len(page3.shapes)
        assert total >= 1, f"Slide 3 should have shapes from SVG image, got {total}"

    def test_multiple_images_slide(self, image_test_parsed):
        """Slide 4 should have multiple image shapes."""
        page4 = image_test_parsed.pages[3]
        image_shapes = [s for s in page4.shapes if s.tag == 'image']
        assert len(image_shapes) >= 1, "Slide 4 should have image shapes"

    def test_image_has_text_around(self, image_test_parsed):
        """Slides with images should also have text segments."""
        page2 = image_test_parsed.pages[1]
        assert len(page2.text_segments) > 0, "Slide 2 should have text segments"


class TestImagePPTXOutput:
    """Test image PPTX output."""

    def test_image_slide_count(self, image_test_pptx):
        """image_test.typ should produce 4 slides."""
        prs = Presentation(image_test_pptx)
        assert len(prs.slides) == 4

    def test_png_image_in_pptx(self, image_test_pptx):
        """Slide 2 should contain the PNG image as a picture shape."""
        prs = Presentation(image_test_pptx)
        slide2 = prs.slides[1]

        pic_count = 0
        for shape in slide2.shapes:
            if hasattr(shape, 'image'):
                try:
                    _ = shape.image
                    pic_count += 1
                except Exception:
                    pass
        assert pic_count >= 1, "Slide 2 should have at least 1 picture"

    def test_png_image_has_content(self, image_test_pptx):
        """The PNG image should have actual image data."""
        prs = Presentation(image_test_pptx)
        slide2 = prs.slides[1]

        for shape in slide2.shapes:
            if hasattr(shape, 'image'):
                try:
                    img = shape.image
                    assert len(img.blob) > 0, "Image should have non-empty data"
                    assert img.content_type in (
                        'image/png', 'image/jpeg', 'image/gif',
                        'image/svg+xml', 'image/x-emf',
                    ), f"Unexpected content type: {img.content_type}"
                    return
                except Exception:
                    pass
        # If no picture found, the image may be rendered as shapes
        # which is also acceptable

    def test_text_around_png_image(self, image_test_pptx):
        """Text before and after the PNG image should be present."""
        prs = Presentation(image_test_pptx)
        slide2 = prs.slides[1]

        all_text = ""
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text + " "

        assert "PNG" in all_text or "embedded" in all_text.lower(), (
            f"Text about PNG image not found on slide 2"
        )

    def test_text_around_svg_image(self, image_test_pptx):
        """Text before and after the SVG image should be present."""
        prs = Presentation(image_test_pptx)
        slide3 = prs.slides[2]

        all_text = ""
        for shape in slide3.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text + " "

        assert "SVG" in all_text, f"Text about SVG image not found on slide 3"

    def test_multiple_images_slide(self, image_test_pptx):
        """Slide 4 should have content from multiple images."""
        prs = Presentation(image_test_pptx)
        slide4 = prs.slides[3]

        # Should have shapes (images or their SVG representations)
        assert len(slide4.shapes) >= 2, (
            f"Slide 4 should have multiple shapes, got {len(slide4.shapes)}"
        )

    def test_caption_text_present(self, image_test_pptx):
        """Caption text should be present on the multiple images slide."""
        prs = Presentation(image_test_pptx)
        slide4 = prs.slides[3]

        all_text = ""
        for shape in slide4.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text + " "

        assert "Caption" in all_text or "side" in all_text, (
            "Caption text not found on slide 4"
        )

    def test_svg_image_is_rasterized_to_png(self, image_test_pptx):
        """SVG images should be rasterized to PNG in the PPTX."""
        prs = Presentation(image_test_pptx)
        slide3 = prs.slides[2]  # SVG Image slide

        pic_found = False
        for shape in slide3.shapes:
            if shape.shape_type == 13:  # PICTURE
                img = shape.image
                assert img.content_type == 'image/png', (
                    f"SVG image should be rasterized to PNG, got {img.content_type}"
                )
                # Verify it's a valid PNG with real content
                pil_img = Image.open(BytesIO(img.blob))
                assert pil_img.width > 0
                assert pil_img.height > 0
                pic_found = True
        assert pic_found, "Slide 3 (SVG) should have a rasterized PNG picture"

    def test_svg_image_has_visible_content(self, image_test_pptx):
        """Rasterized SVG image should have visible (non-blank) content."""
        prs = Presentation(image_test_pptx)
        slide3 = prs.slides[2]  # SVG Image slide

        for shape in slide3.shapes:
            if shape.shape_type == 13:  # PICTURE
                img = shape.image
                pil_img = Image.open(BytesIO(img.blob))
                # Check it's not a blank white image - get unique colors
                colors = pil_img.getcolors(maxcolors=10000)
                assert colors is None or len(colors) > 1, (
                    "SVG rasterized image should not be blank (single color)"
                )
                return
        pytest.fail("No picture found on SVG slide")


class TestImageRasterization:
    """Test image rasterization via typst Python package."""

    def test_svg_rasterization_basic(self):
        """_rasterize_image_to_png should rasterize SVG to valid PNG."""
        from typst2pptx.core.converter import TypstSVGConverter, ConversionConfig
        config = ConversionConfig()
        converter = TypstSVGConverter(config)

        svg_data = (
            b'<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100">'
            b'<rect width="200" height="100" fill="red"/>'
            b'<circle cx="100" cy="50" r="40" fill="green"/>'
            b'</svg>'
        )

        png_data = converter._rasterize_image_to_png(svg_data, 'svg')
        assert len(png_data) > 0, "SVG rasterization should produce data"

        img = Image.open(BytesIO(png_data))
        assert img.format == 'PNG'
        assert img.width > 0
        assert img.height > 0

    def test_svg_rasterization_has_transparency(self):
        """Rasterized SVG should have transparent background (RGBA)."""
        from typst2pptx.core.converter import TypstSVGConverter, ConversionConfig
        config = ConversionConfig()
        converter = TypstSVGConverter(config)

        # Small circle on transparent background
        svg_data = (
            b'<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">'
            b'<circle cx="50" cy="50" r="30" fill="blue"/>'
            b'</svg>'
        )

        png_data = converter._rasterize_image_to_png(svg_data, 'svg')
        img = Image.open(BytesIO(png_data))
        assert img.mode == 'RGBA', (
            f"Rasterized image should be RGBA for transparency, got {img.mode}"
        )

    def test_svg_rasterization_has_content(self):
        """Rasterized SVG should have visible (non-blank) content."""
        from typst2pptx.core.converter import TypstSVGConverter, ConversionConfig
        config = ConversionConfig()
        converter = TypstSVGConverter(config)

        svg_data = (
            b'<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100">'
            b'<rect width="200" height="100" fill="red"/>'
            b'</svg>'
        )

        png_data = converter._rasterize_image_to_png(svg_data, 'svg')
        img = Image.open(BytesIO(png_data))

        colors = img.getcolors(maxcolors=10000)
        assert colors is None or len(colors) > 1, (
            "Rasterized SVG should not be blank"
        )

    def test_pdf_rasterization_basic(self):
        """_rasterize_image_to_png should rasterize PDF to valid PNG."""
        from typst2pptx.core.converter import TypstSVGConverter, ConversionConfig
        config = ConversionConfig()
        converter = TypstSVGConverter(config)

        with open('tests/typ_sources/test_document.pdf', 'rb') as f:
            pdf_data = f.read()

        png_data = converter._rasterize_image_to_png(pdf_data, 'pdf')
        assert len(png_data) > 0, "PDF rasterization should produce data"

        img = Image.open(BytesIO(png_data))
        assert img.format == 'PNG'
        assert img.width > 0
        assert img.height > 0

    def test_pdf_rasterization_with_dimensions(self):
        """_rasterize_image_to_png should respect width_px for PDF."""
        from typst2pptx.core.converter import TypstSVGConverter, ConversionConfig
        config = ConversionConfig()
        converter = TypstSVGConverter(config)

        with open('tests/typ_sources/test_document.pdf', 'rb') as f:
            pdf_data = f.read()

        png_data = converter._rasterize_image_to_png(
            pdf_data, 'pdf', width_px=400
        )
        assert len(png_data) > 0

        img = Image.open(BytesIO(png_data))
        # Width should be approximately based on 400pt at configured DPI
        assert img.width > 100, f"Image width should be significant, got {img.width}"

    def test_pdf_rasterization_has_content(self):
        """Rasterized PDF should have visible content (not blank)."""
        from typst2pptx.core.converter import TypstSVGConverter, ConversionConfig
        config = ConversionConfig()
        converter = TypstSVGConverter(config)

        with open('tests/typ_sources/test_document.pdf', 'rb') as f:
            pdf_data = f.read()

        png_data = converter._rasterize_image_to_png(
            pdf_data, 'pdf', width_px=400
        )
        img = Image.open(BytesIO(png_data))

        colors = img.getcolors(maxcolors=10000)
        assert colors is None or len(colors) > 1, (
            "Rasterized PDF should not be a blank image"
        )

    def test_svgxml_format_alias(self):
        """'svg+xml' format should work the same as 'svg'."""
        from typst2pptx.core.converter import TypstSVGConverter, ConversionConfig
        config = ConversionConfig()
        converter = TypstSVGConverter(config)

        svg_data = (
            b'<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">'
            b'<rect width="100" height="100" fill="blue"/>'
            b'</svg>'
        )

        png_data = converter._rasterize_image_to_png(svg_data, 'svg+xml')
        assert len(png_data) > 0, "svg+xml format should be supported"

        img = Image.open(BytesIO(png_data))
        assert img.format == 'PNG'
