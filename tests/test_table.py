"""Tests for table rendering."""
import pytest
from pptx import Presentation


class TestTableSVGParsing:
    """Test SVG parsing for tables."""

    def test_table_page_count(self, table_test_parsed):
        """table_test.typ should produce 5 slides."""
        assert len(table_test_parsed.pages) == 5

    def test_simple_table_has_text(self, table_test_parsed):
        """Simple table slide should have text segments for cell content."""
        page2 = table_test_parsed.pages[1]
        texts = [s.text for s in page2.text_segments]
        combined = " ".join(texts)
        assert "Name" in combined, f"'Name' not found in: {combined[:200]}"
        assert "Alice" in combined, f"'Alice' not found in: {combined[:200]}"

    def test_simple_table_has_shapes(self, table_test_parsed):
        """Simple table slide should have shape elements (grid lines/borders)."""
        page2 = table_test_parsed.pages[1]
        assert len(page2.shapes) > 0, "Table slide should have shapes for borders/grid"

    def test_colored_table_has_text(self, table_test_parsed):
        """Colored table slide should have cell content."""
        page4 = table_test_parsed.pages[3]
        texts = [s.text for s in page4.text_segments]
        combined = " ".join(texts)
        assert "100k" in combined, f"'100k' not found in: {combined[:200]}"

    def test_colored_table_has_shapes(self, table_test_parsed):
        """Colored table slide should have colored fill shapes."""
        page4 = table_test_parsed.pages[3]
        assert len(page4.shapes) > 0, "Colored table should have fill shapes"

    def test_mixed_table_has_math(self, table_test_parsed):
        """Mixed content table should have math segments."""
        page5 = table_test_parsed.pages[4]
        math_segs = [s for s in page5.text_segments if s.font_variant == 'math']
        assert len(math_segs) > 0, "Mixed table should have math segments"


class TestTablePPTXOutput:
    """Test table PPTX output."""

    def test_table_slide_count(self, table_test_pptx):
        """table_test.typ should produce 5 slides."""
        prs = Presentation(table_test_pptx)
        assert len(prs.slides) == 5

    def test_simple_table_cell_content(self, table_test_pptx):
        """Simple table cells should appear as text on slide 2."""
        prs = Presentation(table_test_pptx)
        slide2 = prs.slides[1]

        all_text = ""
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text + " "

        assert "Name" in all_text, f"'Name' not found in slide text"
        assert "Alice" in all_text, f"'Alice' not found in slide text"
        assert "Bob" in all_text, f"'Bob' not found in slide text"
        assert "Charlie" in all_text, f"'Charlie' not found in slide text"

    def test_simple_table_has_all_cells(self, table_test_pptx):
        """All table cells should be present on simple table slide."""
        prs = Presentation(table_test_pptx)
        slide2 = prs.slides[1]

        all_text = ""
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text + " "

        for expected in ["Name", "Age", "City", "New York", "London", "Tokyo"]:
            assert expected in all_text, f"'{expected}' not found in slide text"

    def test_styled_table_header(self, table_test_pptx):
        """Styled table should have header content."""
        prs = Presentation(table_test_pptx)
        slide3 = prs.slides[2]

        all_text = ""
        for shape in slide3.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text + " "

        assert "Description" in all_text, f"'Description' not found in slide text"
        assert "Status" in all_text, f"'Status' not found in slide text"

    def test_styled_table_bold_header(self, table_test_pptx):
        """Styled table header should have bold formatting."""
        prs = Presentation(table_test_pptx)
        slide3 = prs.slides[2]

        has_bold = False
        for shape in slide3.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.bold and r.text.strip():
                            has_bold = True
        assert has_bold, "Styled table header should have bold text"

    def test_colored_table_has_shapes(self, table_test_pptx):
        """Colored table slide should have shapes (for colored fills/borders)."""
        prs = Presentation(table_test_pptx)
        slide4 = prs.slides[3]
        assert len(slide4.shapes) > 5, (
            f"Colored table should have multiple shapes, got {len(slide4.shapes)}"
        )

    def test_colored_table_cell_content(self, table_test_pptx):
        """Colored table cells should be present."""
        prs = Presentation(table_test_pptx)
        slide4 = prs.slides[3]

        all_text = ""
        for shape in slide4.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text + " "

        for expected in ["Quarter", "Revenue", "100k", "120k"]:
            assert expected in all_text, f"'{expected}' not found in slide text"

    def test_mixed_table_has_math(self, table_test_pptx):
        """Mixed content table should contain math characters."""
        prs = Presentation(table_test_pptx)
        slide5 = prs.slides[4]

        has_math = False
        for shape in slide5.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name == 'Cambria Math':
                            has_math = True
        assert has_math, "Mixed table should have Cambria Math runs"

    def test_mixed_table_text_before_after(self, table_test_pptx):
        """Text before and after the table should be present."""
        prs = Presentation(table_test_pptx)
        slide5 = prs.slides[4]

        all_text = ""
        for shape in slide5.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text + " "

        assert "before" in all_text, "'before' not found in slide text"
        assert "after" in all_text, "'after' not found in slide text"

    def test_mixed_table_colored_text(self, table_test_pptx):
        """Table with colored text should have non-default color runs."""
        prs = Presentation(table_test_pptx)
        slide5 = prs.slides[4]

        has_color = False
        for shape in slide5.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.color and r.font.color.rgb:
                            color_hex = str(r.font.color.rgb)
                            if color_hex not in ('000000', 'FFFFFF', None):
                                has_color = True
        assert has_color, "Table with colored text should have non-black color runs"

    def test_table_shapes_are_filled_rectangles(self, table_test_pptx):
        """Table grid shapes should include filled shapes (cell backgrounds)."""
        from lxml import etree
        prs = Presentation(table_test_pptx)
        ns = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }

        # Slide 4 (colored table) should have filled shapes
        slide4 = prs.slides[3]
        sp_tree = slide4.shapes._spTree
        filled_count = 0
        for sp in sp_tree.findall('.//p:sp', ns):
            solid_fill = sp.find('.//a:solidFill', ns)
            if solid_fill is not None:
                filled_count += 1
        assert filled_count > 0, "Colored table should have filled shapes"
