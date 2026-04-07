"""Tests for inline math handling."""
import pytest
from pptx import Presentation


class TestInlineMathSVGParsing:
    """Test SVG parsing for inline math content."""

    def test_inline_math_page_count(self, inline_math_test_parsed):
        """inline_math_test.typ should produce 5 slides."""
        assert len(inline_math_test_parsed.pages) == 5

    def test_inline_math_has_math_segments(self, inline_math_test_parsed):
        """Pages with inline math should have math-variant segments."""
        # Slide 2 has superscripts/subscripts
        page2 = inline_math_test_parsed.pages[1]
        math_segs = [s for s in page2.text_segments if s.font_variant == 'math']
        assert len(math_segs) > 0, "Slide 2 should have math segments"

    def test_inline_math_has_text_segments(self, inline_math_test_parsed):
        """Pages with inline math should also have regular text segments."""
        page2 = inline_math_test_parsed.pages[1]
        text_segs = [s for s in page2.text_segments if s.font_variant != 'math']
        assert len(text_segs) > 0, "Slide 2 should have non-math text segments"

    def test_inline_math_glyph_data(self, inline_math_test_parsed):
        """Math segments should carry glyph usage data for curve rendering."""
        page2 = inline_math_test_parsed.pages[1]
        math_segs = [s for s in page2.text_segments if s.font_variant == 'math']
        # At least some math segments should have glyph_uses
        has_glyph = any(s.glyph_uses for s in math_segs)
        assert has_glyph, "Math segments should have glyph_uses data"

    def test_inline_math_glyph_scale(self, inline_math_test_parsed):
        """Math segments should have a positive glyph_scale."""
        page2 = inline_math_test_parsed.pages[1]
        math_segs = [s for s in page2.text_segments
                     if s.font_variant == 'math' and s.glyph_uses]
        assert len(math_segs) > 0
        for seg in math_segs:
            assert seg.glyph_scale > 0, (
                f"glyph_scale should be positive, got {seg.glyph_scale}"
            )


class TestInlineMathPPTXOutput:
    """Test inline math PPTX output."""

    def test_inline_math_slide_count(self, inline_math_test_pptx):
        """inline_math_test.typ should produce 5 slides."""
        prs = Presentation(inline_math_test_pptx)
        assert len(prs.slides) == 5

    def test_superscript_merged_with_text(self, inline_math_test_pptx):
        """Superscript like r^2 should be in the same textbox as surrounding text."""
        prs = Presentation(inline_math_test_pptx)
        slide2 = prs.slides[1]

        # Find the textbox containing "The area is" - it should also contain
        # math runs (A, =, pi, r, 2) and "for a circle."
        found = False
        for shape in slide2.shapes:
            if shape.has_text_frame:
                all_text = ""
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text
                if "area" in all_text and "circle" in all_text:
                    # Check that math chars are also in this textbox
                    has_math_char = any(
                        r.font.name == 'Cambria Math'
                        for p in shape.text_frame.paragraphs
                        for r in p.runs
                    )
                    assert has_math_char, (
                        f"Textbox with 'area...circle' should contain Cambria Math runs"
                    )
                    found = True
        assert found, "No textbox found containing both 'area' and 'circle'"

    def test_subscript_merged_with_text(self, inline_math_test_pptx):
        """Subscript like H_2 O should be in the same textbox as surrounding text."""
        prs = Presentation(inline_math_test_pptx)
        slide2 = prs.slides[1]

        found = False
        for shape in slide2.shapes:
            if shape.has_text_frame:
                all_text = ""
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text
                if "Water" in all_text and "𝐻" in all_text:
                    found = True
        assert found, "No textbox found with both 'Water' and H (math)"

    def test_superscript_has_baseline_offset(self, inline_math_test_pptx):
        """Superscript runs should have a positive PPTX baseline attribute (raised text)."""
        prs = Presentation(inline_math_test_pptx)
        slide2 = prs.slides[1]
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        # Find the textbox with "area" and "circle" (contains r^2)
        found_baseline = False
        for shape in slide2.shapes:
            if shape.has_text_frame:
                all_text = ""
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text
                if "area" in all_text and "circle" in all_text:
                    # Check runs for baseline attribute
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            rPr = r._r.find(f'{{{ns_a}}}rPr')
                            if rPr is not None:
                                baseline = rPr.get('baseline')
                                if baseline is not None and int(baseline) > 0:
                                    found_baseline = True
        assert found_baseline, (
            "Superscript in 'r^2' should have positive baseline offset in PPTX (raised text)"
        )

    def test_subscript_has_baseline_offset(self, inline_math_test_pptx):
        """Subscript runs should have a negative PPTX baseline attribute (lowered text)."""
        prs = Presentation(inline_math_test_pptx)
        slide2 = prs.slides[1]
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        # Find the textbox with "Water" and "H" (contains H_2 O)
        found_baseline = False
        for shape in slide2.shapes:
            if shape.has_text_frame:
                all_text = ""
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text
                if "Water" in all_text or "𝐻" in all_text:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            rPr = r._r.find(f'{{{ns_a}}}rPr')
                            if rPr is not None:
                                baseline = rPr.get('baseline')
                                if baseline is not None and int(baseline) < 0:
                                    found_baseline = True
        assert found_baseline, (
            "Subscript in 'H_2 O' should have negative baseline offset in PPTX (lowered text)"
        )

    def test_greek_letters_inline(self, inline_math_test_pptx):
        """Greek letters should appear in a textbox with surrounding text."""
        prs = Presentation(inline_math_test_pptx)
        slide3 = prs.slides[2]

        found = False
        for shape in slide3.shapes:
            if shape.has_text_frame:
                all_text = ""
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text
                if "Greek" in all_text and "𝛼" in all_text:
                    found = True
        assert found, "No textbox found with both 'Greek' and alpha"

    def test_operators_inline(self, inline_math_test_pptx):
        """Math operators should appear with surrounding text."""
        prs = Presentation(inline_math_test_pptx)
        slide3 = prs.slides[2]

        found = False
        for shape in slide3.shapes:
            if shape.has_text_frame:
                all_text = ""
                has_cambria = False
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text
                        if r.font.name == 'Cambria Math':
                            has_cambria = True
                if "Operators" in all_text and has_cambria:
                    found = True
        assert found, "No textbox found with 'Operators' and Cambria Math"

    def test_function_notation(self, inline_math_test_pptx):
        """Function notation f(x) = x^2 + 2x + 1 should render correctly."""
        prs = Presentation(inline_math_test_pptx)
        slide4 = prs.slides[3]

        # Check that function notation appears somewhere on this slide
        all_math_text = []
        for shape in slide4.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name == 'Cambria Math':
                            all_math_text.append(r.text)
        combined = " ".join(all_math_text)
        assert '𝑓' in combined or 'f' in combined, (
            f"Function 'f' not found in math text: '{combined[:100]}'"
        )

    def test_mixed_text_and_math(self, inline_math_test_pptx):
        """Mixed text and math should coexist in textboxes."""
        prs = Presentation(inline_math_test_pptx)
        slide4 = prs.slides[3]

        found = False
        for shape in slide4.shapes:
            if shape.has_text_frame:
                all_text = ""
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text
                if "value" in all_text and "𝑥" in all_text:
                    found = True
        assert found, "No textbox found with both 'value' and math 'x'"

    def test_absolute_value(self, inline_math_test_pptx):
        """Absolute value |x| should be rendered correctly."""
        prs = Presentation(inline_math_test_pptx)
        slide5 = prs.slides[4]

        found = False
        for shape in slide5.shapes:
            if shape.has_text_frame:
                all_text = ""
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text
                if "Absolute" in all_text and "𝑥" in all_text:
                    found = True
        assert found, "No textbox found with 'Absolute' and math 'x'"

    def test_simple_inline_math_uses_text(self, inline_math_test_pptx):
        """Slides with simple inline math (slides 2-3) should use text, not MathGlyph shapes."""
        from lxml import etree
        prs = Presentation(inline_math_test_pptx)
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

        # Slides 2-3 have simple inline math (letters, sub/superscripts)
        # These should NOT have MathGlyph shapes in auto mode
        # Note: Slide 4 has stacked fractions (1/2, 1/3, 5/6) which correctly
        # use MathGlyph shapes, so we only check slides 2-3 here
        for si in [1, 2]:  # 0-indexed: slides 2, 3
            slide = prs.slides[si]
            sp_tree = slide.shapes._spTree
            math_glyph_count = 0
            for sp in sp_tree.findall('.//p:sp', ns):
                cnvpr = sp.find('.//p:nvSpPr/p:cNvPr', ns)
                if cnvpr is not None and 'MathGlyph' in cnvpr.get('name', ''):
                    math_glyph_count += 1
            assert math_glyph_count == 0, (
                f"Slide {si+1} has {math_glyph_count} MathGlyph shapes, "
                f"expected 0 (simple inline math should use text)"
            )

    def test_stacked_fraction_uses_curves(self, inline_math_test_pptx):
        """Slide 4 has stacked fractions (1/2, 1/3, 5/6) which should use MathGlyph in auto mode."""
        from lxml import etree
        prs = Presentation(inline_math_test_pptx)
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

        slide4 = prs.slides[3]
        sp_tree = slide4.shapes._spTree
        math_glyph_count = 0
        for sp in sp_tree.findall('.//p:sp', ns):
            cnvpr = sp.find('.//p:nvSpPr/p:cNvPr', ns)
            if cnvpr is not None and 'MathGlyph' in cnvpr.get('name', ''):
                math_glyph_count += 1
        assert math_glyph_count > 0, (
            "Slide 4 should have MathGlyph shapes for stacked fractions (1/2, 1/3, 5/6)"
        )

    def test_complex_inline_math_uses_curves(self, inline_math_test_pptx):
        """Slide 5 has complex inline math (sum, sqrt) which should use MathGlyph in auto mode."""
        from lxml import etree
        prs = Presentation(inline_math_test_pptx)
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

        slide5 = prs.slides[4]
        sp_tree = slide5.shapes._spTree
        math_glyph_count = 0
        for sp in sp_tree.findall('.//p:sp', ns):
            cnvpr = sp.find('.//p:nvSpPr/p:cNvPr', ns)
            if cnvpr is not None and 'MathGlyph' in cnvpr.get('name', ''):
                math_glyph_count += 1
        assert math_glyph_count > 0, (
            "Slide 5 should have MathGlyph shapes for complex inline math (sum, sqrt)"
        )


class TestDisplayMathAsCurves:
    """Test display math curve rendering."""

    def test_display_math_has_glyph_shapes(self, math_test_pptx):
        """Display math slides should have MathGlyph shapes (curve rendering)."""
        from lxml import etree
        prs = Presentation(math_test_pptx)
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

        # Slide 2 has display integrals/sums
        slide2 = prs.slides[1]
        sp_tree = slide2.shapes._spTree
        math_glyph_count = 0
        for sp in sp_tree.findall('.//p:sp', ns):
            cnvpr = sp.find('.//p:nvSpPr/p:cNvPr', ns)
            if cnvpr is not None and 'MathGlyph' in cnvpr.get('name', ''):
                math_glyph_count += 1
        assert math_glyph_count > 0, (
            "Slide 2 should have MathGlyph shapes for display math"
        )

    def test_display_math_glyphs_have_fill(self, math_test_pptx):
        """MathGlyph shapes should have solid fill (not just outline)."""
        from lxml import etree
        prs = Presentation(math_test_pptx)
        ns = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }

        slide2 = prs.slides[1]
        sp_tree = slide2.shapes._spTree
        filled_count = 0
        for sp in sp_tree.findall('.//p:sp', ns):
            cnvpr = sp.find('.//p:nvSpPr/p:cNvPr', ns)
            if cnvpr is not None and 'MathGlyph' in cnvpr.get('name', ''):
                solid_fill = sp.find('.//a:solidFill', ns)
                if solid_fill is not None:
                    filled_count += 1
        assert filled_count > 0, "MathGlyph shapes should have solid fill"

    def test_inline_math_still_uses_text(self, math_test_pptx):
        """Inline math (slide 1) should still use Cambria Math text runs."""
        prs = Presentation(math_test_pptx)
        slide1 = prs.slides[0]

        has_cambria = False
        for shape in slide1.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name == 'Cambria Math':
                            has_cambria = True
        assert has_cambria, (
            "Slide 1 inline math should use Cambria Math text runs"
        )

    def test_display_math_grouped(self, math_test_pptx):
        """Display math glyphs should be wrapped in MathFormula group shapes."""
        from lxml import etree
        prs = Presentation(math_test_pptx)
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

        # Slide 2 has display integrals/sums
        slide2 = prs.slides[1]
        sp_tree = slide2.shapes._spTree
        grp_count = 0
        for grp in sp_tree.findall(f'.//p:grpSp', ns):
            cnvpr = grp.find(f'p:nvGrpSpPr/p:cNvPr', ns)
            if cnvpr is not None and 'MathFormula' in cnvpr.get('name', ''):
                grp_count += 1
        assert grp_count > 0, (
            "Display math should be grouped in MathFormula group shapes"
        )
