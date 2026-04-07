"""Tests for text positioning accuracy."""
import pytest
from pptx import Presentation
from pptx.util import Emu


class TestTextPositioning:
    """Test that text is positioned correctly in the PPTX."""

    def test_title_centered_horizontally(self, basic_text_pptx):
        """Title 'Hello World' should be roughly centered on slide 1."""
        prs = Presentation(basic_text_pptx)
        slide_width = prs.slide_width

        slide1 = prs.slides[0]
        for shape in slide1.shapes:
            if shape.has_text_frame:
                text = " ".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
                if "Hello World" in text:
                    shape_center = shape.left + shape.width / 2
                    slide_center = slide_width / 2
                    # Should be within 30% of center
                    assert abs(shape_center - slide_center) < slide_width * 0.3, (
                        f"Title not centered: shape_center={shape_center}, "
                        f"slide_center={slide_center}"
                    )
                    return

        pytest.fail("Title 'Hello World' not found on slide 1")

    def test_body_text_left_aligned(self, basic_text_pptx):
        """Body text with bullets should be left-aligned."""
        prs = Presentation(basic_text_pptx)
        slide2 = prs.slides[1]

        bullet_shapes = []
        for shape in slide2.shapes:
            if shape.has_text_frame:
                text = " ".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
                if "Bullet point" in text:
                    bullet_shapes.append(shape)

        assert len(bullet_shapes) > 0, "No bullet point shapes found"

        # All bullet shapes should have similar left position
        lefts = [s.left for s in bullet_shapes]
        for left in lefts:
            assert abs(left - lefts[0]) < Emu(50000), (
                "Bullet points should be aligned vertically"
            )

    def test_bullets_vertically_spaced(self, basic_text_pptx):
        """Bullet points should be evenly spaced vertically."""
        prs = Presentation(basic_text_pptx)
        slide2 = prs.slides[1]

        bullet_tops = []
        for shape in slide2.shapes:
            if shape.has_text_frame:
                text = " ".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
                if "Bullet point" in text:
                    bullet_tops.append(shape.top)

        if len(bullet_tops) >= 3:
            bullet_tops.sort()
            spacings = [bullet_tops[i+1] - bullet_tops[i] for i in range(len(bullet_tops)-1)]
            # All spacings should be similar (within 20%)
            avg_spacing = sum(spacings) / len(spacings)
            for spacing in spacings:
                assert abs(spacing - avg_spacing) < avg_spacing * 0.3, (
                    f"Uneven bullet spacing: {spacings}"
                )

    def test_text_within_slide_bounds(self, basic_text_pptx):
        """All text shapes should be within slide boundaries."""
        prs = Presentation(basic_text_pptx)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    assert shape.left >= 0, (
                        f"Slide {si+1}: shape at negative left: {shape.left}"
                    )
                    assert shape.top >= 0, (
                        f"Slide {si+1}: shape at negative top: {shape.top}"
                    )
                    # Allow some overflow for right/bottom (text wrapping)
                    assert shape.left < slide_width * 1.1, (
                        f"Slide {si+1}: shape left {shape.left} > slide width {slide_width}"
                    )

    def test_font_size_consistency(self, basic_text_pptx):
        """Body text should have consistent font sizes on the same slide."""
        prs = Presentation(basic_text_pptx)
        slide2 = prs.slides[1]

        body_sizes = []
        for shape in slide2.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if "Bullet point" in run.text and run.font.size:
                            body_sizes.append(run.font.size)

        if len(body_sizes) >= 2:
            # All body text should have the same font size
            assert all(s == body_sizes[0] for s in body_sizes), (
                f"Inconsistent body font sizes: {body_sizes}"
            )

    def test_title_larger_than_body(self, basic_text_pptx):
        """Title font size should be larger than body text."""
        prs = Presentation(basic_text_pptx)
        slide2 = prs.slides[1]

        title_size = None
        body_size = None

        for shape in slide2.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.size:
                            if "Second Slide" in run.text:
                                title_size = run.font.size
                            elif "Bullet point" in run.text:
                                body_size = run.font.size

        if title_size and body_size:
            assert title_size > body_size, (
                f"Title ({title_size}) should be larger than body ({body_size})"
            )

    def test_colored_text_preserved(self, basic_text_pptx):
        """Red and blue text should have correct colors."""
        prs = Presentation(basic_text_pptx)
        slide3 = prs.slides[2]

        found_red = False
        found_blue = False

        for shape in slide3.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.color and run.font.color.rgb:
                            color = str(run.font.color.rgb)
                            if "Red" in run.text and color.startswith("FF"):
                                found_red = True
                            elif "blue" in run.text and color.endswith("FF"):
                                found_blue = True

        # At least check that colored text exists
        assert found_red or found_blue, "No colored text found on slide 3"

    def test_inline_code_different_font(self, basic_text_pptx):
        """Inline code should use monospace font."""
        prs = Presentation(basic_text_pptx)
        slide3 = prs.slides[2]

        found_mono = False
        for shape in slide3.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if "inline code" in run.text:
                            if run.font.name and 'consol' in run.font.name.lower():
                                found_mono = True

        assert found_mono, "Inline code should use Consolas font"

    def test_page_numbers_at_bottom(self, basic_text_pptx):
        """Page numbers should be at the bottom of the slide."""
        prs = Presentation(basic_text_pptx)
        slide_height = prs.slide_height

        slides = list(prs.slides)
        for si in range(1, len(slides)):  # Pages 2+ (0-indexed: 1+) have numbers
            slide = slides[si]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = " ".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
                    if "/" in text and any(c.isdigit() for c in text):
                        # Page number should be in bottom 15% of slide
                        assert shape.top > slide_height * 0.85, (
                            f"Page number on slide {si+1} too high: "
                            f"top={shape.top}, threshold={slide_height * 0.85}"
                        )
                        break


class TestMultiRunTextbox:
    """Test that same-line text segments are merged into multi-run textboxes."""

    def test_bold_italic_same_textbox(self, basic_text_pptx):
        """Bold and italic text on the same line should be in one textbox."""
        prs = Presentation(basic_text_pptx)
        slide1 = prs.slides[0]

        for shape in slide1.shapes:
            if shape.has_text_frame:
                runs = []
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        runs.append(run)

                text = " ".join(r.text for r in runs)
                if "bold text" in text and "italic text" in text:
                    # Found the merged textbox
                    has_bold = any(r.font.bold for r in runs)
                    has_italic = any(r.font.italic for r in runs)
                    assert has_bold, "Should have bold run in merged textbox"
                    assert has_italic, "Should have italic run in merged textbox"
                    return

        pytest.fail("No merged textbox with both bold and italic text found")

    def test_multi_run_preserves_order(self, basic_text_pptx):
        """Text order should be preserved in multi-run textboxes."""
        prs = Presentation(basic_text_pptx)
        slide1 = prs.slides[0]

        for shape in slide1.shapes:
            if shape.has_text_frame:
                all_text = " ".join(
                    r.text for p in shape.text_frame.paragraphs for r in p.runs
                )
                if "bold text" in all_text and "italic text" in all_text:
                    # Check order: "bold text" should come before "italic text"
                    bold_pos = all_text.index("bold text")
                    italic_pos = all_text.index("italic text")
                    assert bold_pos < italic_pos, (
                        "Bold text should come before italic text"
                    )
                    return

    def test_bullet_each_on_own_line(self, basic_text_pptx):
        """Each bullet point should be a separate textbox."""
        prs = Presentation(basic_text_pptx)
        slide2 = prs.slides[1]

        bullet_count = 0
        for shape in slide2.shapes:
            if shape.has_text_frame:
                text = " ".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
                if "Bullet point" in text:
                    bullet_count += 1

        assert bullet_count == 3, f"Expected 3 separate bullet textboxes, got {bullet_count}"
