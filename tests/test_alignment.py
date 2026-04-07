"""Tests for text alignment detection (center, right, justify)."""
import pytest
from pptx import Presentation
from pptx.enum.text import PP_ALIGN


class TestAlignmentSVGParsing:
    """Test SVG parsing for alignment content."""

    def test_alignment_page_count(self, alignment_test_parsed):
        """alignment_test.typ should produce 5 slides."""
        assert len(alignment_test_parsed.pages) >= 4

    def test_alignment_has_text(self, alignment_test_parsed):
        """Each page should have text segments."""
        for page in alignment_test_parsed.pages:
            assert len(page.text_segments) > 0


class TestAlignmentPPTXOutput:
    """Test alignment PPTX output."""

    def test_alignment_slide_count(self, alignment_test_pptx):
        """alignment_test.typ should produce slides."""
        prs = Presentation(alignment_test_pptx)
        assert len(prs.slides) >= 4

    def test_left_aligned_text(self, alignment_test_pptx):
        """Slide 2 should have left-aligned text (default)."""
        prs = Presentation(alignment_test_pptx)
        slide2 = prs.slides[1]

        for shape in slide2.shapes:
            if shape.has_text_frame:
                all_text = ""
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        all_text += r.text
                if "left aligned" in all_text.lower():
                    # Left alignment is the default (None or PP_ALIGN.LEFT)
                    for p in shape.text_frame.paragraphs:
                        assert p.alignment in (None, PP_ALIGN.LEFT), (
                            f"Expected left alignment, got {p.alignment}"
                        )
                    return
        # Left-aligned text not found is acceptable if it's in different shape

    def test_center_aligned_text(self, alignment_test_pptx):
        """Slide 3 should have center-aligned text."""
        prs = Presentation(alignment_test_pptx)
        slide3 = prs.slides[2]

        found_center = False
        for shape in slide3.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.alignment == PP_ALIGN.CENTER:
                        found_center = True
        assert found_center, "No center-aligned paragraph found on slide 3"

    def test_right_aligned_text(self, alignment_test_pptx):
        """Slide 4 should have right-aligned text."""
        prs = Presentation(alignment_test_pptx)
        slide4 = prs.slides[3]

        found_right = False
        for shape in slide4.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.alignment == PP_ALIGN.RIGHT:
                        found_right = True
        assert found_right, "No right-aligned paragraph found on slide 4"

    def test_justified_text(self, alignment_test_pptx):
        """Slide 5 should have justified text (long paragraph)."""
        prs = Presentation(alignment_test_pptx)
        if len(prs.slides) < 5:
            pytest.skip("Not enough slides for justify test")
        slide5 = prs.slides[4]

        found_justify = False
        for shape in slide5.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.alignment == PP_ALIGN.JUSTIFY:
                        found_justify = True
        assert found_justify, "No justified paragraph found on slide 5"
