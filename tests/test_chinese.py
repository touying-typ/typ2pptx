"""Tests for Chinese text handling."""
import pytest
from pptx import Presentation


class TestChineseSVGParsing:
    """Test SVG parsing for Chinese content."""

    def test_chinese_page_count(self, chinese_test_parsed):
        """chinese_test.typ should produce 4 slides (title + 3 content)."""
        assert len(chinese_test_parsed.pages) == 4

    def test_chinese_has_text_segments(self, chinese_test_parsed):
        """Each page should have text segments."""
        for page in chinese_test_parsed.pages:
            assert len(page.text_segments) > 0, (
                f"Page {page.page_num} has no text segments"
            )

    def test_chinese_characters_present(self, chinese_test_parsed):
        """Chinese characters should be present in text segments."""
        page1 = chinese_test_parsed.pages[0]
        all_text = "".join(s.text for s in page1.text_segments)
        # Check for Chinese characters from the test file
        assert "中文" in all_text, (
            f"Chinese characters not found in page 1 text: '{all_text[:200]}'"
        )

    def test_chinese_font_variants_detected(self, chinese_test_parsed):
        """Multiple font variants should be detected including CJK text."""
        styles = {v.style for v in chinese_test_parsed.font_variants.values()}
        # CJK text often detected as 'mono' due to quadratic curves in glyph paths
        assert len(styles) >= 2, f"Expected at least 2 font styles, got: {styles}"

    def test_chinese_segments_have_positions(self, chinese_test_parsed):
        """Chinese text segments should have valid positions."""
        for page in chinese_test_parsed.pages:
            for seg in page.text_segments:
                assert isinstance(seg.x, float), f"x is not float: {type(seg.x)}"
                assert isinstance(seg.y, float), f"y is not float: {type(seg.y)}"
                assert seg.font_size > 0, f"Invalid font size: {seg.font_size}"


class TestChinesePPTXOutput:
    """Test the Chinese text PPTX output."""

    def test_chinese_slide_count(self, chinese_test_pptx):
        """chinese_test.typ should produce 4 slides (title + 3 content)."""
        prs = Presentation(chinese_test_pptx)
        assert len(prs.slides) == 4

    def test_chinese_text_in_pptx(self, chinese_test_pptx):
        """Chinese text should appear in the PPTX output."""
        prs = Presentation(chinese_test_pptx)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            all_text.append(run.text)

        combined = " ".join(all_text)
        # Check for Chinese text from the test file
        assert "中文" in combined, (
            f"Chinese characters not found in PPTX: '{combined[:200]}'"
        )

    def test_mixed_chinese_english(self, chinese_test_pptx):
        """Mixed Chinese and English text should both appear."""
        prs = Presentation(chinese_test_pptx)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            all_text.append(run.text)

        combined = " ".join(all_text)
        assert "Hello" in combined, "English text missing from mixed content"
        assert "你好" in combined or "混合" in combined, (
            "Chinese text missing from mixed content"
        )

    def test_chinese_bullet_points(self, chinese_test_pptx):
        """Chinese bullet points should be present."""
        prs = Presentation(chinese_test_pptx)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            all_text.append(run.text)

        combined = " ".join(all_text)
        # Check for numbered items from test file
        assert "第一" in combined or "第二" in combined, (
            f"Chinese bullet point text not found: '{combined[:200]}'"
        )

    def test_chinese_math_formula(self, chinese_test_pptx):
        """Math formula should be present alongside Chinese text."""
        prs = Presentation(chinese_test_pptx)
        has_math = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            if run.font.name == 'Cambria Math':
                                has_math = True
        assert has_math, "No Cambria Math font found (math formulas missing)"

    def test_chinese_title_has_larger_font(self, chinese_test_pptx):
        """Chinese title should have a larger font size than body text."""
        prs = Presentation(chinese_test_pptx)
        slide1 = prs.slides[0]
        title_sizes = []
        body_sizes = []
        for shape in slide1.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.size:
                            if "中文" in run.text or "测试" in run.text:
                                title_sizes.append(run.font.size)
                            elif "typ2pptx" in run.text or "中文支持" in run.text:
                                body_sizes.append(run.font.size)
        if title_sizes and body_sizes:
            assert max(title_sizes) > max(body_sizes), (
                f"Title size {max(title_sizes)} should be > body size {max(body_sizes)}"
            )
