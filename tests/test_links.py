"""Tests for hyperlink handling."""
import pytest
from pptx import Presentation


class TestLinkSVGParsing:
    """Test SVG parsing for link content."""

    def test_link_page_count(self, link_test_parsed):
        """link_test.typ should produce slides."""
        assert len(link_test_parsed.pages) >= 1

    def test_link_regions_detected(self, link_test_parsed):
        """SVG parser should detect link regions from <a> elements."""
        # The link test has links on the first content page
        has_links = False
        for page in link_test_parsed.pages:
            if page.links:
                has_links = True
        assert has_links, "No link regions detected in SVG"

    def test_link_has_valid_href(self, link_test_parsed):
        """Link regions should have non-empty href."""
        for page in link_test_parsed.pages:
            for link in page.links:
                assert link.href, "Link should have non-empty href"

    def test_link_has_valid_bbox(self, link_test_parsed):
        """Link regions should have positive width and height."""
        for page in link_test_parsed.pages:
            for link in page.links:
                assert link.width > 0, f"Link width should be positive, got {link.width}"
                assert link.height > 0, f"Link height should be positive, got {link.height}"


class TestLinkPPTXOutput:
    """Test link PPTX output."""

    def test_link_pptx_has_hyperlinks(self, link_test_pptx):
        """PPTX output should have hyperlink relationships."""
        prs = Presentation(link_test_pptx)
        has_hyperlink = False
        for slide in prs.slides:
            for rid, rel in slide.part.rels.items():
                if 'hyperlink' in str(rel.reltype):
                    has_hyperlink = True
        assert has_hyperlink, "No hyperlink relationships found in PPTX"

    def test_typst_website_link(self, link_test_pptx):
        """The Typst website link should be present."""
        prs = Presentation(link_test_pptx)
        found_typst = False
        for slide in prs.slides:
            for rid, rel in slide.part.rels.items():
                if 'hyperlink' in str(rel.reltype) and 'typst.app' in rel.target_ref:
                    found_typst = True
        assert found_typst, "Typst website hyperlink not found"

    def test_mailto_link(self, link_test_pptx):
        """The mailto link should be present."""
        prs = Presentation(link_test_pptx)
        found_mailto = False
        for slide in prs.slides:
            for rid, rel in slide.part.rels.items():
                if 'hyperlink' in str(rel.reltype) and 'mailto:' in rel.target_ref:
                    found_mailto = True
        assert found_mailto, "Mailto hyperlink not found"

    def test_internal_link_not_hyperlinked(self, link_test_pptx):
        """Internal document links (#page-N) should NOT be PPTX hyperlinks."""
        prs = Presentation(link_test_pptx)
        for slide in prs.slides:
            for rid, rel in slide.part.rels.items():
                if 'hyperlink' in str(rel.reltype):
                    # Should not have fragment-style internal links
                    assert not rel.target_ref.startswith('#'), (
                        f"Internal link should not be a hyperlink: {rel.target_ref}"
                    )

    def test_link_text_preserves_original_color(self, link_test_pptx):
        """Linked text should preserve original color and suppress PPTX defaults."""
        from lxml import etree
        prs = Presentation(link_test_pptx)
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            rPr = r._r.find(f'{{{ns_a}}}rPr')
                            if rPr is not None:
                                hlink = rPr.find(f'{{{ns_a}}}hlinkClick')
                                if hlink is not None:
                                    # Should explicitly suppress PPTX default underline
                                    assert rPr.get('u') == 'none', (
                                        "Link should have u='none' to suppress default underline"
                                    )
                                    # Should have explicit solidFill to prevent theme color
                                    solid_fill = rPr.find(f'{{{ns_a}}}solidFill')
                                    assert solid_fill is not None, (
                                        "Link should have explicit solidFill to prevent "
                                        "theme hyperlink color override"
                                    )

    def test_theme_hyperlink_colors_neutralized(self, link_test_pptx):
        """Theme-level hyperlink colors should be overridden to black.

        PowerPoint themes define hlink (blue) and folHlink (purple) colors
        that override run-level solidFill. These must be neutralized to
        prevent the default blue/purple link styling.
        """
        from lxml import etree
        prs = Presentation(link_test_pptx)
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        for master in prs.slide_masters:
            for rel in master.part.rels.values():
                if 'theme' in str(rel.reltype):
                    theme_xml = etree.fromstring(rel.target_part.blob)

                    hlink = theme_xml.find(f'.//{{{ns_a}}}hlink')
                    assert hlink is not None, "Theme should have hlink element"
                    hlink_color = hlink.find(f'{{{ns_a}}}srgbClr')
                    assert hlink_color is not None, "hlink should have srgbClr"
                    assert hlink_color.get('val') == '000000', (
                        f"Theme hlink color should be 000000 (black), "
                        f"got {hlink_color.get('val')}"
                    )

                    folhlink = theme_xml.find(f'.//{{{ns_a}}}folHlink')
                    assert folhlink is not None, "Theme should have folHlink element"
                    folhlink_color = folhlink.find(f'{{{ns_a}}}srgbClr')
                    assert folhlink_color is not None, "folHlink should have srgbClr"
                    assert folhlink_color.get('val') == '000000', (
                        f"Theme folHlink color should be 000000 (black), "
                        f"got {folhlink_color.get('val')}"
                    )

    def test_correct_text_is_linked(self, link_test_pptx):
        """The correct text runs should have hyperlinks (not surrounding text)."""
        from lxml import etree
        prs = Presentation(link_test_pptx)
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        linked_texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            rPr = r._r.find(f'{{{ns_a}}}rPr')
                            if rPr is not None:
                                hlink = rPr.find(f'{{{ns_a}}}hlinkClick')
                                if hlink is not None:
                                    linked_texts.append(r.text)

        # "Typst website" should be linked, NOT "Visit " or "for more info"
        assert any('Typst' in t or 'website' in t for t in linked_texts), (
            f"Expected 'Typst website' in linked texts, got: {linked_texts}"
        )
        # "test@example.com" should be linked, NOT "Email:"
        assert any('example.com' in t for t in linked_texts), (
            f"Expected 'test@example.com' in linked texts, got: {linked_texts}"
        )
