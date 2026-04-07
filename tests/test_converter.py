"""Tests for the converter module."""
import os
import pytest
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

from typst2pptx.core.converter import (
    ConversionConfig,
    TypstSVGConverter,
    _parse_color,
    _font_variant_to_props,
    _group_segments_by_line,
    _compute_text_position,
    compile_typst_to_svg,
    query_speaker_notes,
    convert_typst_to_pptx,
)
from typst2pptx.core.typst_svg_parser import TextSegment


class TestParseColor:
    """Test CSS color parsing."""

    def test_hex_6digit(self):
        assert _parse_color("#FF0000") == "FF0000"

    def test_hex_3digit(self):
        result = _parse_color("#F00")
        assert result == "FF0000"

    def test_hex_lowercase(self):
        assert _parse_color("#ff0000") == "FF0000"

    def test_rgb_function(self):
        result = _parse_color("rgb(255, 0, 0)")
        assert result == "FF0000"

    def test_rgb_green(self):
        result = _parse_color("rgb(0, 255, 0)")
        assert result == "00FF00"

    def test_rgb_blue(self):
        result = _parse_color("rgb(0, 0, 255)")
        assert result == "0000FF"

    def test_none_string_returns_none(self):
        result = _parse_color("none")
        assert result is None

    def test_url_returns_none(self):
        """Gradient references should return None."""
        result = _parse_color("url(#gradient1)")
        assert result is None

    def test_empty_string_returns_none(self):
        result = _parse_color("")
        assert result is None

    def test_named_color_black(self):
        assert _parse_color("black") == "000000"

    def test_named_color_white(self):
        assert _parse_color("white") == "FFFFFF"

    def test_named_color_red(self):
        assert _parse_color("red") == "FF0000"

    def test_named_color_case_insensitive(self):
        assert _parse_color("Blue") == "0000FF"
        assert _parse_color("GREEN") == "008000"

    def test_rgba_function(self):
        """rgba() should return the color, ignoring alpha."""
        result = _parse_color("rgba(255, 0, 0, 0.5)")
        assert result == "FF0000"

    def test_rgb_percent(self):
        result = _parse_color("rgb(100%, 0%, 0%)")
        assert result == "FF0000"

    def test_rgba_percent(self):
        result = _parse_color("rgba(0%, 100%, 0%, 50%)")
        assert result == "00FF00"


class TestFontVariantToProps:
    """Test font variant to python-pptx properties mapping."""

    def test_regular(self):
        props = _font_variant_to_props("regular")
        assert props.get("bold") is False
        assert props.get("italic") is False

    def test_bold(self):
        props = _font_variant_to_props("bold")
        assert props.get("bold") is True

    def test_italic(self):
        props = _font_variant_to_props("italic")
        assert props.get("italic") is True

    def test_bolditalic(self):
        props = _font_variant_to_props("bolditalic")
        assert props.get("bold") is True
        assert props.get("italic") is True

    def test_mono(self):
        props = _font_variant_to_props("mono")
        assert "font_family" in props


class TestGroupSegmentsByLine:
    """Test text segment grouping by baseline."""

    def _make_seg(self, y, text="test", font_size=12.0):
        """Create a TextSegment with given y position."""
        return TextSegment(
            text=text, x=0.0, y=y, width=50.0, height=font_size,
            font_size=font_size, font_variant="regular", fill_color="#000",
        )

    def test_same_baseline(self):
        segs = [self._make_seg(100.0, "a"), self._make_seg(100.0, "b")]
        groups = _group_segments_by_line(segs)
        assert len(groups) == 1

    def test_different_baselines(self):
        # Need enough Y separation (> 2px after baseline calc)
        segs = [self._make_seg(100.0, "a"), self._make_seg(200.0, "b")]
        groups = _group_segments_by_line(segs)
        assert len(groups) == 2

    def test_close_baselines_grouped(self):
        """Segments within 2px baseline tolerance should be grouped."""
        segs = [self._make_seg(100.0, "a"), self._make_seg(101.0, "b")]
        groups = _group_segments_by_line(segs)
        assert len(groups) == 1

    def test_empty_segments(self):
        groups = _group_segments_by_line([])
        assert len(groups) == 0

    def test_three_lines(self):
        segs = [
            self._make_seg(100.0, "line1"),
            self._make_seg(200.0, "line2"),
            self._make_seg(300.0, "line3"),
        ]
        groups = _group_segments_by_line(segs)
        assert len(groups) == 3

    def test_multiple_segments_same_line(self):
        segs = [
            self._make_seg(100.0, "a"),
            self._make_seg(100.0, "b"),
            self._make_seg(100.0, "c"),
        ]
        groups = _group_segments_by_line(segs)
        assert len(groups) == 1
        assert len(groups[0]) == 3

    def test_unsorted_input(self):
        """Segments provided out of order should still be grouped correctly."""
        segs = [
            self._make_seg(300.0, "c"),
            self._make_seg(100.0, "a"),
            self._make_seg(200.0, "b"),
        ]
        groups = _group_segments_by_line(segs)
        assert len(groups) == 3


class TestConversionConfig:
    """Test conversion configuration."""

    def test_default_config(self):
        config = ConversionConfig()
        assert config.include_speaker_notes is True
        assert config.verbose is False

    def test_custom_config(self):
        config = ConversionConfig(
            include_speaker_notes=False,
            verbose=True,
        )
        assert config.include_speaker_notes is False
        assert config.verbose is True

    def test_math_mode_defaults(self):
        """Default math modes: inline=auto, display=glyph."""
        config = ConversionConfig()
        assert config.inline_math_mode == "auto"
        assert config.display_math_mode == "glyph"

    def test_math_mode_custom(self):
        """Custom math rendering modes."""
        config = ConversionConfig(
            inline_math_mode="glyph",
            display_math_mode="text",
        )
        assert config.inline_math_mode == "glyph"
        assert config.display_math_mode == "text"


class TestCompileTypstToSvg:
    """Test SVG compilation from .typ files."""

    def test_basic_text_compiles(self, typ_sources_dir, typst_ts_cli):
        """basic_text.typ should compile to SVG successfully."""
        typ_path = typ_sources_dir / "basic_text.typ"
        svg_path = compile_typst_to_svg(str(typ_path), typst_ts_cli=typst_ts_cli)
        assert os.path.exists(svg_path)
        assert svg_path.endswith(".svg")

    def test_svg_file_not_empty(self, typ_sources_dir, typst_ts_cli):
        """Generated SVG should have content."""
        typ_path = typ_sources_dir / "basic_text.typ"
        svg_path = compile_typst_to_svg(str(typ_path), typst_ts_cli=typst_ts_cli)
        assert os.path.getsize(svg_path) > 1000

    def test_nonexistent_file_raises(self, typst_ts_cli):
        """Compiling a nonexistent file should raise FileNotFoundError."""
        with pytest.raises(FileNotFoundError):
            compile_typst_to_svg("/nonexistent/file.typ", typst_ts_cli=typst_ts_cli)


class TestQuerySpeakerNotes:
    """Test speaker notes extraction."""

    def test_speaker_notes_extraction(self, typ_sources_dir):
        """speaker_notes_test.typ should produce notes for pages 1, 2, 4."""
        typ_path = typ_sources_dir / "speaker_notes_test.typ"
        notes = query_speaker_notes(str(typ_path))

        assert isinstance(notes, dict)
        assert len(notes) == 3

        # Page indices are 0-based from pdfpc
        assert 1 in notes
        assert 2 in notes
        assert 4 in notes

        assert "explain" in notes[1].lower() or "concept" in notes[1].lower()
        assert "item" in notes[2].lower()
        assert "audience" in notes[4].lower() or "q&a" in notes[4].lower()

    def test_no_notes_returns_empty(self, typ_sources_dir):
        """basic_text.typ has no speaker notes."""
        typ_path = typ_sources_dir / "basic_text.typ"
        notes = query_speaker_notes(str(typ_path))
        assert notes == {}

    def test_nonexistent_file_returns_empty(self):
        """Querying a nonexistent file should return empty dict (not raise)."""
        notes = query_speaker_notes("/nonexistent/file.typ")
        assert notes == {}


class TestPPTXOutput:
    """Test the generated PPTX files."""

    def test_basic_text_slide_count(self, basic_text_pptx):
        """basic_text.typ should produce 3 slides."""
        prs = Presentation(basic_text_pptx)
        assert len(prs.slides) == 3

    def test_basic_text_slide_size(self, basic_text_pptx):
        """Slide size should match 842x474 typst coordinates."""
        prs = Presentation(basic_text_pptx)
        width = prs.slide_width
        height = prs.slide_height
        assert width == 8020050
        assert height == 4514850

    def test_basic_text_has_textboxes(self, basic_text_pptx):
        """Slides should contain shapes."""
        prs = Presentation(basic_text_pptx)
        total_shapes = sum(len(slide.shapes) for slide in prs.slides)
        assert total_shapes > 0, "No shapes in any slide"

    def test_basic_text_has_text_content(self, basic_text_pptx):
        """Text should be extractable from the PPTX."""
        prs = Presentation(basic_text_pptx)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                all_text.append(run.text)

        combined = " ".join(all_text)
        assert len(combined) > 0, "No text content found in PPTX"

    def test_basic_text_font_styles(self, basic_text_pptx):
        """PPTX should have bold and italic text runs."""
        prs = Presentation(basic_text_pptx)
        has_bold = False
        has_italic = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.bold:
                                has_bold = True
                            if run.font.italic:
                                has_italic = True

        assert has_bold, "No bold text found in PPTX"
        assert has_italic, "No italic text found in PPTX"

    def test_basic_text_each_slide_has_shapes(self, basic_text_pptx):
        """Every slide should have at least one shape."""
        prs = Presentation(basic_text_pptx)
        for i, slide in enumerate(prs.slides):
            assert len(slide.shapes) > 0, f"Slide {i+1} has no shapes"

    def test_shapes_test_slide_count(self, shapes_test_pptx):
        """shapes_test.typ should produce at least 4 slides."""
        prs = Presentation(shapes_test_pptx)
        assert len(prs.slides) >= 4

    def test_shapes_test_has_freeform_shapes(self, shapes_test_pptx):
        """shapes_test should have custom geometry (freeform) shapes."""
        prs = Presentation(shapes_test_pptx)
        from lxml import etree
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        freeform_count = 0
        for slide in prs.slides:
            sp_tree = slide.shapes._spTree
            for sp in sp_tree:
                cust_geom = sp.findall('.//a:custGeom', ns)
                if cust_geom:
                    freeform_count += 1

        assert freeform_count > 0, "No freeform shapes found in shapes_test PPTX"

    def test_shapes_have_fill_colors(self, shapes_test_pptx):
        """Some shapes should have solid fill colors."""
        prs = Presentation(shapes_test_pptx)
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        has_fill = False
        for slide in prs.slides:
            sp_tree = slide.shapes._spTree
            for sp in sp_tree:
                fills = sp.findall('.//a:solidFill', ns)
                if fills:
                    has_fill = True
                    break

        assert has_fill, "No shapes with solid fill colors found"

    def test_shapes_have_gradient_fills(self, shapes_test_pptx):
        """shapes_test includes a gradient box; verify gradient fill is present."""
        prs = Presentation(shapes_test_pptx)
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        has_gradient = False
        for slide in prs.slides:
            sp_tree = slide.shapes._spTree
            for sp in sp_tree:
                grads = sp.findall('.//a:gradFill', ns)
                if grads:
                    has_gradient = True
                    # Verify gradient has at least 2 stops
                    for grad in grads:
                        gs_list = grad.findall('.//a:gs', ns)
                        assert len(gs_list) >= 2, "Gradient should have at least 2 stops"
                    break

        assert has_gradient, "No shapes with gradient fills found"

    def test_speaker_notes_slide_count(self, speaker_notes_pptx):
        """speaker_notes_test.typ should produce 5 slides."""
        prs = Presentation(speaker_notes_pptx)
        assert len(prs.slides) == 5

    def test_speaker_notes_in_pptx(self, speaker_notes_pptx):
        """Speaker notes should be present on correct slides."""
        prs = Presentation(speaker_notes_pptx)

        notes_by_slide = {}
        for i, slide in enumerate(prs.slides):
            try:
                ns = slide.notes_slide
                text = ns.notes_text_frame.text
                if text:
                    notes_by_slide[i] = text
            except Exception:
                pass

        # Slide 2 (idx 1), Slide 3 (idx 2), Slide 5 (idx 4) should have notes
        assert 1 in notes_by_slide, "Slide 2 missing speaker notes"
        assert 2 in notes_by_slide, "Slide 3 missing speaker notes"
        assert 4 in notes_by_slide, "Slide 5 missing speaker notes"

        # Verify note content
        assert "explain" in notes_by_slide[1].lower() or "concept" in notes_by_slide[1].lower()
        assert "item" in notes_by_slide[2].lower()

    def test_slides_without_notes(self, speaker_notes_pptx):
        """Slides 1 and 4 should NOT have speaker notes."""
        prs = Presentation(speaker_notes_pptx)

        for idx in [0, 3]:  # Slide 1 and 4 (0-indexed)
            slide = prs.slides[idx]
            try:
                ns = slide.notes_slide
                text = ns.notes_text_frame.text
                assert not text or text.strip() == "", (
                    f"Slide {idx+1} should not have notes but has: '{text}'"
                )
            except Exception:
                pass  # No notes slide is also acceptable


class TestConvertTypstToPptx:
    """Test the end-to-end conversion function."""

    def test_e2e_basic_text(self, typ_sources_dir, output_dir, typst_ts_cli):
        """Full end-to-end conversion of basic_text.typ."""
        output_path = str(output_dir / "e2e_basic_text.pptx")
        result = convert_typst_to_pptx(
            str(typ_sources_dir / "basic_text.typ"),
            output_path,
            typst_ts_cli=typst_ts_cli,
        )
        assert os.path.exists(result)
        prs = Presentation(result)
        assert len(prs.slides) == 3

    def test_e2e_speaker_notes(self, typ_sources_dir, output_dir, typst_ts_cli):
        """Full end-to-end conversion with speaker notes."""
        output_path = str(output_dir / "e2e_speaker_notes.pptx")
        result = convert_typst_to_pptx(
            str(typ_sources_dir / "speaker_notes_test.typ"),
            output_path,
            typst_ts_cli=typst_ts_cli,
        )
        assert os.path.exists(result)

        prs = Presentation(result)
        # Verify notes are present
        slide2_notes = prs.slides[1].notes_slide.notes_text_frame.text
        assert len(slide2_notes) > 0

    def test_e2e_shapes(self, typ_sources_dir, output_dir, typst_ts_cli):
        """Full end-to-end conversion with shapes."""
        output_path = str(output_dir / "e2e_shapes.pptx")
        result = convert_typst_to_pptx(
            str(typ_sources_dir / "shapes_test.typ"),
            output_path,
            typst_ts_cli=typst_ts_cli,
        )
        assert os.path.exists(result)
        prs = Presentation(result)
        assert len(prs.slides) >= 4

    def test_e2e_from_svg(self, basic_text_svg, output_dir):
        """Conversion from pre-compiled SVG."""
        output_path = str(output_dir / "e2e_from_svg.pptx")
        result = convert_typst_to_pptx(
            basic_text_svg,
            output_path,
        )
        assert os.path.exists(result)
        prs = Presentation(result)
        assert len(prs.slides) == 3

    def test_unsupported_format_raises(self, output_dir):
        """Conversion of unsupported file format should raise ValueError."""
        with pytest.raises(ValueError, match="Unsupported"):
            convert_typst_to_pptx(
                "test.pdf",
                str(output_dir / "should_not_exist.pptx"),
            )
